"""
helpers.py — lightweight layout & drawing helpers atop python-pptx.
"""
from functools import lru_cache
from typing import Dict, Tuple, Optional
from xml.etree import ElementTree as ET

import yaml
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER

EMU_PER_INCH = 914400

def emu(x_inches: float) -> int:
    return int(x_inches * EMU_PER_INCH)

def place_norm(slide, x: float, y: float, w: float, h: float):
    """Place a textbox using normalized (0..1) coords relative to slide size."""
    presentation = slide.part.package.presentation_part.presentation
    sw, sh = presentation.slide_width, presentation.slide_height
    left, top, width, height = int(sw*x), int(sh*y), int(sw*w), int(sh*h)
    return slide.shapes.add_textbox(left, top, width, height)

def set_paragraph(para, text: str, face: str, size_pt: Optional[int], bold: bool=False, color_rgb=(38,38,38),
                  align="left", inherit_fonts: bool=False):
    para.text = text
    run = para.runs[0] if para.runs else para.add_run()
    run.text = text
    if not inherit_fonts and face:
        run.font.name = face
    if not inherit_fonts and size_pt is not None:
        run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color_rgb)
    align_map = {
        "left": PP_PARAGRAPH_ALIGNMENT.LEFT,
        "center": PP_PARAGRAPH_ALIGNMENT.CENTER,
        "right": PP_PARAGRAPH_ALIGNMENT.RIGHT,
    }
    para.alignment = align_map.get(align, PP_PARAGRAPH_ALIGNMENT.LEFT)

def add_big_number(slide, value_text: str, label_text: str, fonts, colors, x=0.05, y=0.15):
    box = place_norm(slide, x, y, 0.27, 0.25)
    tf = box.text_frame
    tf.clear()
    p = _ensure_paragraph(tf)
    set_paragraph(p, value_text, fonts.title_face, fonts.title_size_pt, True, _resolve_color(colors, "accent1"))
    p2 = tf.add_paragraph()
    set_paragraph(p2, label_text, fonts.body_face, fonts.body_size_pt, False, _resolve_color(colors, "text_dark"))

def add_picture_fit(slide, image_path: str, x: float, y: float, w: float, h: float):
    presentation = slide.part.package.presentation_part.presentation
    sw, sh = presentation.slide_width, presentation.slide_height
    left, top, width, height = int(sw*x), int(sh*y), int(sw*w), int(sh*h)
    return slide.shapes.add_picture(image_path, left, top, width=width, height=height)

def find_placeholder(slide, placeholder_type: PP_PLACEHOLDER, fallback_to_first_text: bool=False):
    for ph in slide.placeholders:
        if ph.placeholder_format.type == placeholder_type:
            return ph
    if fallback_to_first_text:
        for ph in slide.placeholders:
            if getattr(ph, "has_text_frame", False) and ph.placeholder_format.type != PP_PLACEHOLDER.TITLE:
                return ph
    return None

def _resolve_color(colors, key: str):
    if hasattr(colors, key):
        return getattr(colors, key)
    if isinstance(colors, dict) and key in colors:
        return colors[key]
    raise AttributeError(f"Color '{key}' not found in provided color scheme")

def reset_text_frame(tf):
    tf.clear()
    return _ensure_paragraph(tf)

def _ensure_paragraph(tf):
    return tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()


@lru_cache(maxsize=4)
def _load_theme_colors(zip_blob: bytes) -> Dict[str, str]:
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    colors: Dict[str, str] = {}
    root = ET.fromstring(zip_blob)
    clr_scheme = root.find(".//a:clrScheme", ns)
    if clr_scheme is None:
        return colors
    for child in clr_scheme:
        name = child.attrib.get("name") or child.tag.split("}")[-1]
        srgb = child.find(".//a:srgbClr", ns)
        if srgb is not None and "val" in srgb.attrib:
            colors[name] = srgb.attrib["val"].upper()
            continue
        scheme_clr = child.find(".//a:schemeClr", ns)
        if scheme_clr is not None and "val" in scheme_clr.attrib:
            colors[name] = colors.get(scheme_clr.attrib["val"], "").upper()
    return colors


def get_theme_colors(prs) -> Dict[str, str]:
    cached = getattr(prs, "_theme_color_cache", None)
    if cached is not None:
        return cached
    theme_blob = None
    for part in prs.part.package.iter_parts():
        if part.partname.endswith("/theme/theme1.xml"):
            theme_blob = part.blob
            break
    colors = _load_theme_colors(theme_blob) if theme_blob else {}
    setattr(prs, "_theme_color_cache", colors)
    return colors


def load_theme_colors(path: str) -> Dict[str, RGBColor]:
    with open(path, "r") as f:
        raw = yaml.safe_load(f) or {}

    def to_rgb(value):
        if isinstance(value, RGBColor):
            return value
        if isinstance(value, str) and value.startswith("#") and len(value) == 7:
            try:
                r = int(value[1:3], 16)
                g = int(value[3:5], 16)
                b = int(value[5:7], 16)
                return RGBColor(r, g, b)
            except ValueError:
                return value
        return value

    return {k: to_rgb(v) for k, v in raw.items()}
