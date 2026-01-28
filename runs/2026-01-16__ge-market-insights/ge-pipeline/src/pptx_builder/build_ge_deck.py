"""
build_ge_deck.py — Builder for the Market Insights Pilot Validation deck (manufacturer-agnostic).
"""
import sys
from pathlib import Path

# Allow running this file directly (without needing PYTHONPATH=.)
_REPO_ROOT = Path(__file__).resolve().parents[2]
if str(_REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(_REPO_ROOT))

import argparse
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import subprocess
import math
import zipfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from src.pptx_builder.style_tokens import DeckTokens, infer_theme_from_template, FontTokens
from src.pptx_builder.slides import create_title_slide, create_section_header
from src.pptx_builder.layouts import find_layout

# Hardcoded layout mapping based on config/deck.yaml and standard Premier template
LAYOUT_MAP = {
    "TITLE": ["Premier Title Dark", "Title Slide"],
    "TITLE_CONTENT": ["Premier Two Content Subtitle 1a", "Title and Content"],
    "MAIN_CONTENT": ["Premier Main Content 1a", "Title and Content"],  # Single full-width body
    "SECTION": ["Premier Title Slide 2 - Dark", "Section Header"],
    "BLANK": ["Blank"]
}

def _setup_tokens(template_path: Path) -> DeckTokens:
    theme = infer_theme_from_template(template_path)
    return DeckTokens(
        colors=theme,
        fonts=FontTokens(),
        layout_map=LAYOUT_MAP,
        force_title_contrast=True
    )

def remove_all_slides(prs):
    """Remove all slides from the presentation using the nuclear option."""
    # NUCLEAR OPTION: Delete all existing slides immediately
    for i in range(len(prs.slides)-1, -1, -1): 
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]


def _convert_template_to_pptx(src_path: Path, dst_path: Path) -> Path:
    """Convert a .potx-like file into a .pptx by adjusting the main content type."""
    old_type = "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"
    new_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"

    dst_path.parent.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(src_path, "r") as zin, zipfile.ZipFile(dst_path, "w", compression=zipfile.ZIP_DEFLATED) as zout:
        for info in zin.infolist():
            data = zin.read(info.filename)
            if info.filename == "[Content_Types].xml":
                try:
                    text = data.decode("utf-8")
                    text = text.replace(old_type, new_type)
                    data = text.encode("utf-8")
                except Exception:
                    # If decoding fails, write original bytes.
                    pass
            zout.writestr(info, data)
    return dst_path


def _ensure_pptx_template(template_path: Path, snapshot_dir: Path) -> Path:
    """Return a pptx-compatible template path usable by python-pptx."""
    suffix = template_path.suffix.lower()
    if suffix == ".pptx":
        return template_path

    # Common case: branded template is distributed as .potx.
    if suffix == ".potx":
        out = snapshot_dir / f"{template_path.stem}.pptx"
        print(f"Converting {template_path} to {out} (content-type rewrite)...")
        return _convert_template_to_pptx(template_path, out)

    return template_path

def generate_market_trend_charts(csv_path: Path, output_dir: Path):
    df = pd.read_csv(csv_path)

    if 'report_category' in df.columns:
        df = df.rename(columns={
            'report_category': 'Category',
            'year_quarter': 'Year_Quarter',
            'share_of_observed_spend': 'Market_Share_Pct',
            'manufacturer_normalized': 'Manufacturer'
        })

    sns.set_theme(style="whitegrid")
    output_paths = {}

    # Keep manufacturer colors consistent across all Market Trend slides.
    manufacturer_order = [
        "GE",
        "SIEMENS",
        "PHILIPS",
        "CANON",
        "SAMSUNG",
        "NIHON KOHDEN",
        "SPACELABS",
        "MINDRAY",
        "OTHER",
    ]
    manufacturer_colors = {
        # Use a stable, readable palette (Tab10-ish) so colors never shift slide-to-slide.
        "GE": "#1f77b4",
        "SIEMENS": "#ff7f0e",
        "PHILIPS": "#2ca02c",
        "CANON": "#d62728",
        "SAMSUNG": "#9467bd",
        "NIHON KOHDEN": "#8c564b",
        "SPACELABS": "#e377c2",
        "MINDRAY": "#bcbd22",
        "OTHER": "#7f7f7f",
    }

    if "Manufacturer" in df.columns:
        df["Manufacturer"] = df["Manufacturer"].astype(str).str.strip().str.upper()

    # Assign deterministic colors for any unexpected manufacturers.
    known = set(manufacturer_colors)
    all_mfrs = sorted(set(df["Manufacturer"].dropna())) if "Manufacturer" in df.columns else []
    unknown_mfrs = [m for m in all_mfrs if m not in known]
    if unknown_mfrs:
        extra_palette = sns.color_palette("tab20", n_colors=len(unknown_mfrs)).as_hex()
        for mfr, color in zip(unknown_mfrs, extra_palette):
            manufacturer_colors[mfr] = color
        manufacturer_order.extend(unknown_mfrs)

    def _slugify(value: str) -> str:
        return value.lower().replace(" ", "_")

    for category, category_df in df.groupby('Category'):
        fig, ax = plt.subplots(figsize=(10, 5))
        present_order = [m for m in manufacturer_order if m in set(category_df["Manufacturer"].dropna())]
        sns.lineplot(
            data=category_df,
            x="Year_Quarter",
            y="Market_Share_Pct",
            hue="Manufacturer",
            hue_order=present_order,
            palette=manufacturer_colors,
            marker="o",
            ax=ax
        )
        ax.set_title(f"{category} Market Share (%)")
        ax.set_xlabel("Quarter")
        ax.set_ylabel("Market Share (%)")
        ax.legend(loc='upper left', bbox_to_anchor=(1.02, 1))
        plt.setp(ax.get_xticklabels(), rotation=45, ha='right')
        plt.tight_layout()
        chart_path = output_dir / f"market_trend_{_slugify(category)}.png"
        fig.savefig(chart_path, dpi=300, bbox_inches='tight')
        plt.close(fig)
        output_paths[category] = chart_path

    return output_paths

def _set_title(slide, text: str, font_size_pt: int = 32):
    """Helper to set slide title and enforce a smaller font size."""
    if not slide.shapes.title:
        return
    slide.shapes.title.text = text
    if slide.shapes.title.text_frame:
        for p in slide.shapes.title.text_frame.paragraphs:
            p.font.size = Pt(font_size_pt)

def add_bullet_slide(prs, tokens, title, headline, bullets, bullet_font_size=18, body_offset_top=0):
    # Use MAIN_CONTENT layout - has single full-width body placeholder
    preferred = LAYOUT_MAP["MAIN_CONTENT"]
    layout = find_layout(prs, ["Main", "Content"], 1, preferred_names=preferred, layout_key="MAIN_CONTENT")
    slide = prs.slides.add_slide(layout)
    
    # Title
    _set_title(slide, title)
    
    # Find the body placeholder (BODY type, not the title)
    body = None
    layout_body = None
    
    # Find layout body placeholder to get correct position/size
    from pptx.enum.shapes import PP_PLACEHOLDER
    for ph in layout.placeholders:
        if hasattr(ph.placeholder_format, 'type') and ph.placeholder_format.type == PP_PLACEHOLDER.BODY:
            layout_body = ph
            break
    
    # Find slide body placeholder
    for shape in slide.placeholders:
        if shape.has_text_frame and shape != slide.shapes.title:
            if hasattr(shape.placeholder_format, 'type'):
                if shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
                    body = shape
                    break
            if body is None:
                body = shape
            
    if body:
        # Ensure body placeholder inherits correct position from layout
        if layout_body:
            body.left = layout_body.left
            body.top = layout_body.top
            body.width = layout_body.width
            body.height = layout_body.height
        
        # Clear the body first
        tf = body.text_frame
        tf.clear()
        tf.word_wrap = True
        
        # Add headline as first paragraph (bold, slightly larger)
        if headline:
            p_headline = tf.paragraphs[0]
            p_headline.text = headline
            p_headline.font.bold = True
            p_headline.font.size = Pt(20)
            p_headline.font.color.rgb = RGBColor(30, 93, 170)  # Accent color
            p_headline.space_after = Pt(16)
            start_idx = 1
        else:
            start_idx = 0
        
        # Add bullet points
        for idx, bullet in enumerate(bullets):
            if start_idx == 0 and idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(bullet_font_size)
            p.space_before = Pt(6)
            p.space_after = Pt(10)

def add_table_slide(prs, tokens, title, headline, csv_path, callout=None):
    # Use MAIN_CONTENT layout for proper title positioning
    preferred = LAYOUT_MAP["MAIN_CONTENT"]
    layout = find_layout(prs, ["Main", "Content"], 1, preferred_names=preferred, layout_key="MAIN_CONTENT")
    slide = prs.slides.add_slide(layout)
    
    _set_title(slide, title)

    # Load Data
    df_raw = pd.read_csv(csv_path)
    
    # Force Aggregation to 3 Rows
    df = df_raw.groupby('Category', as_index=False).agg({
        'Transaction_Spend': 'sum',
        'Supplier_Spend': 'sum'
    })
    df['Coverage_Ratio'] = df['Transaction_Spend'] / df['Supplier_Spend']
    
    # Reorder columns to match requirement
    df = df[["Category", "Transaction_Spend", "Supplier_Spend", "Coverage_Ratio"]]
    
    # Create Table
    rows, cols = df.shape
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(12.0)
    height = Inches(3.3)
    
    table_shape = slide.shapes.add_table(rows + 1, cols, left, top, width, height)
    table = table_shape.table
    table.rows[0].height = Inches(0.4)
    
    # Header
    for i, col_name in enumerate(df.columns):
        cell = table.cell(0, i)
        cell.text = col_name
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        
    # Rows
    for i, (idx, row) in enumerate(df.iterrows()):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            # Format numbers
            if isinstance(val, (int, float)):
                if "Ratio" in df.columns[j] or "Pct" in df.columns[j]:
                    cell.text = f"{val:.1%}" if val < 10 else f"{val:.1f}"
                elif "Spend" in df.columns[j]:
                    cell.text = f"${val:,.0f}"
                else:
                    cell.text = str(val)
            else:
                cell.text = str(val)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(10)
            para.space_after = Pt(0)
        table.rows[i + 1].height = Inches(0.32)

    # Callout
    if callout:
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(5.4), Inches(12.0), Inches(0.8))
        p = textbox.text_frame.add_paragraph()
        p.text = callout
        p.font.bold = True
        p.font.color.rgb = RGBColor(30, 93, 170) # Accent 1 approx
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER

def add_ct_deep_dive_slide(prs, tokens, title, headline, csv_path, insight):
    # Use MAIN_CONTENT layout for proper title positioning
    preferred = LAYOUT_MAP["MAIN_CONTENT"]
    layout = find_layout(prs, ["Main", "Content"], 1, preferred_names=preferred, layout_key="MAIN_CONTENT")
    slide = prs.slides.add_slide(layout)
    
    _set_title(slide, title)
        
    df = pd.read_csv(csv_path)

    # Derive average unit price to identify capital equipment
    df['avg_unit_price'] = df['total_observed_spend'] / df['transaction_count']

    # Exclude UNKNOWN and accessory/parts terms
    exclude_terms = [
        'UNKNOWN', 'CABLE', 'BATTERY', 'PROBE', 'SENSOR', 'KIT', 'TUBE', 'FILTER', 'GRID',
        'PHANTOM', 'COIL', 'INJECTOR', 'CONNECTOR', 'MEDIA', 'CONTRAST', 'ADAPTER', 'MOUNT',
        'BRACKET', 'COVER', 'LICENSE', 'SOFTWARE', 'SERVICE', 'REPAIR', 'MAINTENANCE',
        'LEADWIRE', 'ELECTRODE', 'BELT', 'STRAP', 'PAD', 'CUSHION', 'SYRINGE', 'NEEDLE',
        'CASSETTE', 'BOARD', 'PCB', 'ASSEMBLY', 'CIRCUIT', 'WIRING', 'MODULE', 'SUBSYSTEM'
    ]
    df = df[~df['Product_Description'].str.contains('|'.join(exclude_terms), case=False, na=False)]

    if 'manufacturer_normalized' in df.columns:
        df = df.rename(columns={'manufacturer_normalized': 'Manufacturer', 'total_observed_spend': 'Total_Spend'})

    df = df.dropna(subset=['Manufacturer', 'Product_Description', 'Total_Spend'])

    # Filter for high unit price items (capital equipment typically > $50k)
    df = df[df['avg_unit_price'] >= 50000].copy()

    allowed_manufacturers = ['GE', 'SIEMENS', 'PHILIPS', 'CANON']
    df = df[df['Manufacturer'].str.upper().isin(allowed_manufacturers)].copy()

    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(12.0)
    height = Inches(4.2)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame

    manufacturer_order = [m for m in allowed_manufacturers if m in df['Manufacturer'].str.upper().unique()]
    
    # Show top 3 items per manufacturer, sorted by avg unit price (highest first)
    max_items = 3

    for mfr_key in manufacturer_order:
        mfr_data = df[df['Manufacturer'].str.upper() == mfr_key]
        if mfr_data.empty:
            continue
        display_name = mfr_data['Manufacturer'].iloc[0]
        header = tf.add_paragraph()
        header.text = display_name
        header.font.bold = True
        header.font.size = Pt(13)
        header.font.color.rgb = RGBColor(0, 0, 0)
        header.space_after = Pt(2) 

        # Rank by avg unit price (highest first = most likely capital equipment)
        mfr_sorted = mfr_data.sort_values(by='avg_unit_price', ascending=False).head(max_items)
        for _, row in mfr_sorted.iterrows():
            p = tf.add_paragraph()
            desc = row['Product_Description']
            # Truncate long descriptions for readability
            if len(desc) > 90:
                desc = desc[:87] + "..."
            p.text = f"• {desc}"
            p.level = 1
            p.font.size = Pt(10)
            p.space_after = Pt(0)
            
        # Spacer paragraph between manufacturers
        spacer = tf.add_paragraph()
        spacer.text = ""
        spacer.font.size = Pt(4)

    # Insight textbox moved to open whitespace
    if insight:
        # Move insight box slightly higher to avoid overlap if list gets long
        insight_box = slide.shapes.add_textbox(Inches(8.2), Inches(1.8), Inches(4.0), Inches(1.2))
        tf_insight = insight_box.text_frame
        tf_insight.clear()
        para = tf_insight.paragraphs[0]
        para.text = f"Insight: {insight}"
        para.font.bold = True
        para.font.italic = True
        para.font.size = Pt(14)
        para.font.color.rgb = RGBColor(30, 93, 170)
        para.alignment = PP_ALIGN.RIGHT

    # Source/methodology footer
    source_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12.0), Inches(0.5))
    tf_source = source_box.text_frame
    tf_source.clear()
    src_para = tf_source.paragraphs[0]
    src_para.text = "Source: Premier provider-reported transaction data (Oct 2023 – Sep 2025). Examples of capital equipment found in the category."
    src_para.font.size = Pt(9)
    src_para.font.italic = True
    src_para.font.color.rgb = RGBColor(100, 100, 100)


def add_monitoring_deep_dive_slide(prs, tokens, title, headline, csv_path, insight):
    """
    Add Monitoring (Physiological Monitoring Systems) capital equipment examples slide.
    
    Selection: High avg unit price items (>$25k) that represent capital systems
    (monitors, central stations) rather than accessories.
    """
    # Use MAIN_CONTENT layout for proper title positioning
    preferred = LAYOUT_MAP["MAIN_CONTENT"]
    layout = find_layout(prs, ["Main", "Content"], 1, preferred_names=preferred, layout_key="MAIN_CONTENT")
    slide = prs.slides.add_slide(layout)
    
    _set_title(slide, title)
        
    df = pd.read_csv(csv_path)

    # Derive average unit price to identify capital equipment
    df['avg_unit_price'] = df['total_observed_spend'] / df['transaction_count']

    # Exclude UNKNOWN and accessory/parts terms
    exclude_terms = [
        'UNKNOWN', 'CABLE', 'BATTERY', 'PROBE', 'SENSOR', 'KIT', 'FILTER',
        'CONNECTOR', 'ADAPTER', 'MOUNT', 'BRACKET', 'COVER', 'LICENSE', 'SOFTWARE',
        'SERVICE', 'REPAIR', 'MAINTENANCE', 'LEADWIRE', 'ELECTRODE', 'BELT',
        'STRAP', 'PAD', 'CUSHION', 'SYRINGE', 'NEEDLE', 'CASSETTE', 'BOARD', 'PCB',
        'CIRCUIT', 'WIRING', 'HOSE', 'RECEIVER', 'RENTAL', 'FEE', 'FRAME', 'MODULE'
    ]
    df = df[~df['Product_Description'].str.contains('|'.join(exclude_terms), case=False, na=False)]

    if 'manufacturer_normalized' in df.columns:
        df = df.rename(columns={'manufacturer_normalized': 'Manufacturer', 'total_observed_spend': 'Total_Spend'})

    df = df.dropna(subset=['Manufacturer', 'Product_Description', 'Total_Spend'])

    # Filter for high unit price items (capital monitoring equipment typically > $25k)
    df = df[df['avg_unit_price'] >= 25000].copy()

    # For Monitoring, focus on major manufacturers
    allowed_manufacturers = ['GE', 'PHILIPS', 'NIHON KOHDEN', 'SPACELABS']
    df = df[df['Manufacturer'].str.upper().isin(allowed_manufacturers)].copy()

    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(12.0)
    height = Inches(4.2)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame

    manufacturer_order = [m for m in allowed_manufacturers if m in df['Manufacturer'].str.upper().unique()]
    
    # Show top 3 items per manufacturer, sorted by avg unit price
    max_items = 3

    for mfr_key in manufacturer_order:
        mfr_data = df[df['Manufacturer'].str.upper() == mfr_key]
        if mfr_data.empty:
            continue
        display_name = mfr_data['Manufacturer'].iloc[0]
        header = tf.add_paragraph()
        header.text = display_name
        header.font.bold = True
        header.font.size = Pt(13)
        header.font.color.rgb = RGBColor(0, 0, 0)
        header.space_after = Pt(2)

        # Rank by avg unit price (highest first = most likely capital equipment)
        mfr_sorted = mfr_data.sort_values(by='avg_unit_price', ascending=False).head(max_items)
        for _, row in mfr_sorted.iterrows():
            p = tf.add_paragraph()
            desc = row['Product_Description']
            # Truncate long descriptions for readability
            if len(desc) > 90:
                desc = desc[:87] + "..."
            p.text = f"• {desc}"
            p.level = 1
            p.font.size = Pt(10)
            p.space_after = Pt(0)
            
        # Spacer paragraph between manufacturers
        spacer = tf.add_paragraph()
        spacer.text = ""
        spacer.font.size = Pt(4)

    # Insight textbox
    if insight:
        insight_box = slide.shapes.add_textbox(Inches(8.2), Inches(1.8), Inches(4.0), Inches(1.2))
        tf_insight = insight_box.text_frame
        tf_insight.clear()
        para = tf_insight.paragraphs[0]
        para.text = f"Insight: {insight}"
        para.font.bold = True
        para.font.italic = True
        para.font.size = Pt(14)
        para.font.color.rgb = RGBColor(30, 93, 170)
        para.alignment = PP_ALIGN.RIGHT

    # Source/methodology footer
    source_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12.0), Inches(0.5))
    tf_source = source_box.text_frame
    tf_source.clear()
    src_para = tf_source.paragraphs[0]
    src_para.text = "Source: Premier provider-reported transaction data (Oct 2023 – Sep 2025). Examples of capital equipment found in the category."
    src_para.font.size = Pt(9)
    src_para.font.italic = True
    src_para.font.color.rgb = RGBColor(100, 100, 100)


def add_ct_charity_presence_slide(prs, tokens, title, headline, csv_path: Path):
    """Add Charity CT product presence + hierarchy availability slide."""
    if not csv_path.exists():
        print(f"Warning: Charity CT presence CSV not found: {csv_path} (skipping slide)")
        return

    # Use MAIN_CONTENT layout for proper title positioning
    preferred = LAYOUT_MAP["MAIN_CONTENT"]
    layout = find_layout(prs, ["Main", "Content"], 1, preferred_names=preferred, layout_key="MAIN_CONTENT")

    df = pd.read_csv(csv_path)
    if df.empty:
        print(f"Warning: Charity CT presence CSV had 0 rows: {csv_path} (skipping slide)")
        return

    # Expected columns (NEXT_ITERATION_PLAN.md)
    # - term_label
    # - match_count_anywhere, base_spend_anywhere, last_month_anywhere
    # - match_count_ct_mapped, base_spend_ct_mapped, pct_spend_ct_mapped
    # - top_contract_category_anywhere
    required_cols = [
        "term_label",
        "match_count_anywhere",
        "base_spend_anywhere",
        "last_month_anywhere",
        "match_count_ct_mapped",
        "base_spend_ct_mapped",
        "pct_spend_ct_mapped",
        "top_contract_category_anywhere",
    ]
    for c in required_cols:
        if c not in df.columns:
            df[c] = None

    display_df = df.copy()
    # Preserve configured order when present
    if "term_order" in display_df.columns:
        display_df["term_order"] = pd.to_numeric(display_df["term_order"], errors="coerce")
        display_df = display_df.sort_values("term_order", kind="stable")
    else:
        display_df = display_df.sort_values("term_label", kind="stable")

    # Explicitly label "Frontier ES" as "Frontier (unspecified)" if present
    if "term_key" in display_df.columns:
        mask = display_df["term_key"] == "REVOLUTION_FRONTIER_ES"
        if mask.any():
            display_df.loc[mask, "term_label"] = "Frontier (unspecified)"

    # Rename to slide-friendly headers
    display_df = display_df.rename(
        columns={
            "term_label": "Product line",
            "match_count_anywhere": "# Txns (anywhere)",
            "base_spend_anywhere": "$ Spend (anywhere)",
            "last_month_anywhere": "Last month (anywhere)",
            "match_count_ct_mapped": "# Txns (CT)",
            "base_spend_ct_mapped": "$ Spend (CT)",
            "pct_spend_ct_mapped": "% Spend (CT)",
            "top_contract_category_anywhere": "Top Contract_Category (anywhere)",
        }
    )

    table_cols = [
        "Product line",
        "# Txns (anywhere)",
        "$ Spend (anywhere)",
        "Last month (anywhere)",
        "# Txns (CT)",
        "$ Spend (CT)",
        "% Spend (CT)",
        "Top Contract_Category (anywhere)",
    ]
    display_df = display_df[table_cols]

    def _render_table_on_slide(slide, df_chunk: pd.DataFrame, title_text: str, headline_text: str):
        _set_title(slide, title_text)

        subtitle = None
        for shape in slide.placeholders:
            if shape.name.lower().find("subtitle") >= 0:
                subtitle = shape
                break
        if subtitle:
            subtitle.text = headline_text

        # Remove the layout's content placeholder so no template text shows behind the table.
        # (The table + footer we add below fully replaces the body content area.)
        for shape in list(slide.shapes):
            if not shape.is_placeholder:
                continue
            if shape == slide.shapes.title:
                continue
            if subtitle is not None and shape == subtitle:
                continue
            if shape.has_text_frame:
                try:
                    shape.text_frame.clear()
                except Exception:
                    pass
            try:
                shape.element.getparent().remove(shape.element)
            except Exception:
                # Fallback: move it off-slide if XML removal fails.
                try:
                    shape.left = Inches(20)
                except Exception:
                    pass

        left = Inches(0.6)
        top = Inches(1.65)  # Below title area (title can be 2 lines, needs ~1.5" clearance)
        width = Inches(12.0)
        height = Inches(2.8)  # Slightly shorter to fit more content

        table_shape = slide.shapes.add_table(len(df_chunk) + 1, len(table_cols), left, top, width, height)
        table = table_shape.table
        
        # Optimize column widths to prevent wrapping in text columns (1 and 8)
        # Total available: 12.0 inches
        # Col 1 (Product Line): 2.5"
        # Col 8 (Contract Category): 3.5"
        # Remaining 6 cols share 6.0" -> 1.0" each
        col_widths = [2.5, 1.0, 1.0, 1.0, 1.0, 1.0, 1.0, 3.5] 
        total_w = sum(col_widths) # 12.0
        
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

        table.rows[0].height = Inches(0.36)

        for i, col_name in enumerate(table_cols):
            cell = table.cell(0, i)
            cell.text = col_name
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(10)

        for i, (_, row) in enumerate(df_chunk.iterrows()):
            for j, col_name in enumerate(table_cols):
                val = row[col_name]
                cell = table.cell(i + 1, j)

                def _is_missing(value) -> bool:
                    if value is None:
                        return True
                    try:
                        # pandas NaN/NaT
                        if pd.isna(value):
                            return True
                    except Exception:
                        pass
                    if isinstance(value, float) and math.isnan(value):
                        return True
                    return False

                if col_name in ["$ Spend (anywhere)", "$ Spend (CT)"]:
                    amount = 0.0 if _is_missing(val) else float(val)
                    cell.text = f"${amount:,.0f}"
                elif col_name in ["# Txns (anywhere)", "# Txns (CT)"]:
                    count = 0 if _is_missing(val) else int(val)
                    cell.text = f"{count:,}"
                elif col_name == "% Spend (CT)":
                    try:
                        pct = 0.0 if _is_missing(val) else float(val)
                    except Exception:
                        pct = 0.0
                    if 0 < pct < 0.001:
                        cell.text = "<0.1%"
                    else:
                        cell.text = f"{pct * 100:.1f}%"
                else:
                    cell.text = "—" if _is_missing(val) else str(val)
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(10)
                para.space_after = Pt(0)
            table.rows[i + 1].height = Inches(0.26)

        # Shifted footer to 6.2 (safe gap from table, above slide footer)
        footer = slide.shapes.add_textbox(Inches(0.6), Inches(6.2), Inches(12.0), Inches(0.75))
        ft = footer.text_frame
        ft.clear()
        p1 = ft.paragraphs[0]
        p1.text = (
            "Matches use facility-submitted description/vendor/manufacturer + catalog fields to reduce false negatives for capital equipment. "
            "Terms with 0 matches found no evidence in this manufacturer-filtered data slice."
        )
        p1.font.size = Pt(10)
        p1.font.color.rgb = RGBColor(80, 80, 80)
        p1.alignment = PP_ALIGN.CENTER

        p2 = ft.add_paragraph()
        p2.text = "CT-mapped includes Premier CT mapping and inferred CT for CT-like product term hits."
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(80, 80, 80)
        p2.alignment = PP_ALIGN.CENTER

    # Split into two slides so the table stays readable and never runs into the Premier footer/logo.
    # The Charity term list is fixed at 18 rows, so use a stable 9/9 split.
    if len(display_df) > 10:
        chunks = [display_df.iloc[:9].reset_index(drop=True), display_df.iloc[9:].reset_index(drop=True)]
        for idx, chunk in enumerate(chunks, start=1):
            slide = prs.slides.add_slide(layout)
            _render_table_on_slide(slide, chunk, f"{title} ({idx}/2)", headline)
    else:
        slide = prs.slides.add_slide(layout)
        _render_table_on_slide(slide, display_df.reset_index(drop=True), title, headline)

def add_chart_slide(prs, tokens, title, headline, chart_path):
    # Use MAIN_CONTENT layout for simpler structure
    preferred = LAYOUT_MAP["MAIN_CONTENT"]
    layout = find_layout(prs, ["Main", "Content"], 1, preferred_names=preferred, layout_key="MAIN_CONTENT")
    slide = prs.slides.add_slide(layout)
    
    # Set title
    _set_title(slide, title)
    
    # Insert Chart - position below title area
    # Title ends around 1.2", add some margin
    width = Inches(10.5)
    height = Inches(5.0)  # Slightly taller to use available space
    
    # Calculate centered left position on 13.33" wide slide
    slide_width_inches = 13.333
    left = Inches((slide_width_inches - 10.5) / 2)
    top = Inches(1.4)  # Below title (title ends ~1.2") with small margin
    
    slide.shapes.add_picture(str(chart_path), left, top, width=width, height=height)


def add_q4_timing_overlay_slide(prs, tokens, chart_path: Path):
    """Add a two-column slide: chart left, narrative notes right."""
    if not chart_path.exists():
        print(f"Warning: Q4 timing overlay chart not found: {chart_path} (skipping slide)")
        return

    # Use MAIN_CONTENT layout for proper title positioning
    preferred = LAYOUT_MAP["MAIN_CONTENT"]
    layout = find_layout(prs, ["Main", "Content"], 1, preferred_names=preferred, layout_key="MAIN_CONTENT")
    slide = prs.slides.add_slide(layout)

    _set_title(slide, "Seasonality comparison (indexed): Supplier vs Provider-reported spend")

    # Find subtitle placeholder if present
    subtitle = None
    for shape in slide.placeholders:
        if shape.name.lower().find("subtitle") >= 0:
            subtitle = shape
            break

    # Remove the layout's content placeholder so our custom layout has full control.
    for shape in list(slide.shapes):
        if not shape.is_placeholder:
            continue
        if shape == slide.shapes.title:
            continue
        if subtitle is not None and shape == subtitle:
            continue
        if shape.has_text_frame:
            try:
                shape.text_frame.clear()
            except Exception:
                pass
        try:
            shape.element.getparent().remove(shape.element)
        except Exception:
            try:
                shape.left = Inches(20)
            except Exception:
                pass

    # Left: chart image
    chart_left = Inches(0.6)
    chart_top = Inches(1.55)
    chart_width = Inches(8.35)
    chart_height = Inches(5.05)
    slide.shapes.add_picture(str(chart_path), chart_left, chart_top, width=chart_width, height=chart_height)

    # Right: commentary
    notes_left = Inches(9.15)
    notes_top = Inches(1.65)
    notes_width = Inches(4.0)
    notes_height = Inches(5.0)
    textbox = slide.shapes.add_textbox(notes_left, notes_top, notes_width, notes_height)
    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True

    header = tf.paragraphs[0]
    header.text = "What to notice"
    header.font.bold = True
    header.font.size = Pt(14)
    header.font.color.rgb = RGBColor(30, 93, 170)
    header.space_after = Pt(10)

    bullets = [
        "Indexed patterns are broadly aligned across ~2 years, supporting consistent underlying seasonality.",
        "The largest difference is a Q4→Q1 timing shift in the provider-reported view (e.g., 2024 Q4 into 2025 Q1).",
        "Interpretation: two valid timing conventions (supplier invoice-aligned vs facility PO-date-aligned) can move spend across adjacent quarters without changing longer-run totals.",
    ]

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(12)
        p.space_after = Pt(8)


def add_methodology_slide(prs, tokens):
    """Add a methodology slide explaining how Coverage Ratios are calculated."""
    preferred = LAYOUT_MAP["MAIN_CONTENT"]
    layout = find_layout(prs, ["Main", "Content"], 1, preferred_names=preferred, layout_key="MAIN_CONTENT")
    slide = prs.slides.add_slide(layout)

    _set_title(slide, "Methodology: Coverage Ratio Calculation")

    # Main content textbox
    left = Inches(0.6)
    top = Inches(1.5)
    width = Inches(12.0)
    height = Inches(5.0)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True

    sections = [
        ("Definition", [
            "Coverage Ratio = (Transaction Data Spend) / (Supplier-Reported Spend)",
            "A ratio of 100% indicates the two data sources report identical total spend.",
            "Ratios above 100% indicate the Transaction Data captures additional spend not present in Supplier Tracings.",
        ]),
        ("Data Sources", [
            "Transaction Data: Purchase Orders and Invoices submitted by health systems to Premier.",
            "Supplier-Reported Spend: Sales tracings provided by manufacturers, typically reflecting contracted/GPO sales.",
        ]),
        ("Alignment Rules", [
            "Facility Alignment: Only facilities reporting in both data sources are included (\"dark facilities\" excluded).",
            "Time Window: Rolling 24-month analysis period (consistent across both sources).",
            "Category Mapping: Transactions mapped to capital equipment categories (CT, MRI, Monitoring, Ultrasound HC) using text-mining on product descriptions with $25K+ price thresholds.",
        ]),
        ("Interpretation", [
            ">100%: Transaction Data captures off-contract purchases, bundled accessories, or local deals not reflected in manufacturer tracings.",
            "~100%: Strong alignment between provider spend and manufacturer-reported sales.",
            "<100%: Potential data gaps in Transaction Data (uncommon in this analysis).",
        ]),
    ]

    first = True
    for section_title, bullets in sections:
        # Section header
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = section_title
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(30, 93, 170)
        p.space_before = Pt(10) if section_title != "Definition" else Pt(0)
        p.space_after = Pt(4)

        # Bullets
        for bullet in bullets:
            bp = tf.add_paragraph()
            bp.text = f"• {bullet}"
            bp.font.size = Pt(11)
            bp.space_after = Pt(2)


def main():
    parser = argparse.ArgumentParser(description="Build GE Pilot Validation Deck")
    parser.add_argument("--template", required=True, type=Path)
    parser.add_argument("--snapshot-dir", required=True, type=Path)
    parser.add_argument("--out", required=True, type=Path)
    parser.add_argument(
        "--q4-overlay-chart",
        type=Path,
        default=None,
        help="Optional path to q4_quarterly_overlay_indexed.png; defaults to <snapshot-dir>/visuals/q4_quarterly_overlay_indexed.png",
    )
    args = parser.parse_args()
    
    print(f"Building deck from {args.snapshot_dir} using {args.template}...")

    template_path = _ensure_pptx_template(args.template, args.snapshot_dir)
    try:
        prs = Presentation(template_path)
    except ValueError as e:
        # Some environments copy a .potx to .pptx; python-pptx rejects the template content type.
        if "presentationml.template.main+xml" in str(e):
            repaired = args.snapshot_dir / "template_base_converted.pptx"
            print(f"Template appears to be a .potx-in-disguise; converting to {repaired}...")
            template_path = _convert_template_to_pptx(template_path, repaired)
            prs = Presentation(template_path)
        else:
            raise
    remove_all_slides(prs)
    tokens = _setup_tokens(args.template)

    # Compute parity ratios once so Exec Summary and Data Validation are consistent.
    parity_csv = args.snapshot_dir / "slide_data_parity_validation.csv"
    parity_ratios = {}
    try:
        df_parity = pd.read_csv(parity_csv)
        if not df_parity.empty:
            dfp = df_parity.groupby("Category", as_index=False).agg(
                {"Transaction_Spend": "sum", "Supplier_Spend": "sum"}
            )
            dfp["Coverage_Ratio"] = dfp["Transaction_Spend"] / dfp["Supplier_Spend"]
            parity_ratios = {str(r["Category"]): float(r["Coverage_Ratio"]) for _, r in dfp.iterrows()}
    except Exception:
        parity_ratios = {}

    def _ratio_pct(category: str) -> str:
        ratio = parity_ratios.get(category)
        if ratio is None or (isinstance(ratio, float) and math.isnan(ratio)):
            return "—"
        return f"{ratio * 100:.1f}%"
    
    # Slide 1: Title
    create_title_slide(
        prs, 
        tokens,
        "Market Insights: Data Validation for Capital-Intensive Categories", 
        "Coverage Analysis: MRI, CT, Patient Monitoring, Ultrasound (Hand-Carried)\nJanuary 2026"
    )
    
    # Slide 2: Executive Summary
    add_bullet_slide(
        prs, tokens,
        "Executive Summary",
        "Why We Validated & What We Found",
        [
            "Purpose: Compare provider-reported Transaction Data (POs/Invoices submitted by health systems) against manufacturer-reported Supplier Tracings to validate coverage of capital equipment and accessories.",
            f"Key Finding – Full Market Visibility: Coverage ratios meet or exceed 100% across all four categories (Computed Tomography: {_ratio_pct('Computed Tomography')}; Magnetic Resonance Imaging: {_ratio_pct('Magnetic Resonance Imaging')}; Physiological Monitoring Systems: {_ratio_pct('Physiological Monitoring Systems')}; Ultrasound Hand-Carried: {_ratio_pct('Ultrasound Hand-Carried')}).",
            "Implication – Capturing Off-Contract Spend: Ratios above 100% indicate the Transaction Data captures additional spend (e.g., accessories, local purchases, bundled items) that manufacturer tracings do not report."
        ],
        bullet_font_size=16,
        body_offset_top=Inches(0.4)
    )
    
    # Slide 3: Data Validation
    add_table_slide(
        prs, tokens,
        "Data Validation",
        "Transaction vs. Supplier Reported Spend",
        args.snapshot_dir / "slide_data_parity_validation.csv",
        "High parity validates that our Transaction Data is a reliable proxy for total market activity, encompassing both capital systems and associated accessories often missing from sales tracings."
    )
    
    # Slide 4: CT Deep Dive
    add_ct_deep_dive_slide(
        prs, tokens,
        "Computed Tomography – Capital Equipment Examples",
        "Deep Dive: Capital Equipment in CT",
        args.snapshot_dir / "slide_data_ct_breakout.csv",
        "Transaction data includes capital systems—not just accessories or consumables."
    )

    # Slide 5: Monitoring Deep Dive (mirrors CT structure)
    add_monitoring_deep_dive_slide(
        prs, tokens,
        "Physiological Monitoring Systems – Capital Equipment Examples",
        "Deep Dive: Capital Equipment in Monitoring",
        args.snapshot_dir / "slide_data_monitoring_breakout.csv",
        "Transaction data includes capital systems—not just accessories or consumables."
    )

    # Slide 6: Charity CT presence
    add_ct_charity_presence_slide(
        prs,
        tokens,
        "CT product presence & hierarchy (Charity request)",
        "Presence-first scan (match on description + manufacturer/vendor/brand); CT category not required",
        args.snapshot_dir / "ct_charity_presence_summary.csv",
    )
    
    # Market Trend Slides (CT, MRI, Monitoring, Ultrasound HC)
    chart_paths = generate_market_trend_charts(args.snapshot_dir / "slide_data_market_trends.csv", args.snapshot_dir)
    category_titles = {
        "CT": "Market Trends - Computed Tomography (CT)",
        "MRI": "Market Trends - Magnetic Resonance Imaging (MRI)",
        "Monitoring": "Market Trends - Patient Monitoring Systems",
        "Ultrasound HC": "Market Trends - Ultrasound (Hand-Carried)"
    }
    for category in ["CT", "MRI", "Monitoring", "Ultrasound HC"]:
        chart_path = chart_paths.get(category)
        if not chart_path:
            continue
        add_chart_slide(
            prs, tokens,
            category_titles.get(category, f"Market Trends - {category}"),
            "Quarterly Market Share Trends (Trailing 2 Years)",
            chart_path
        )

    # Q4 timing reconciliation overlay (Supplier vs Provider indexed seasonality)
    q4_overlay_chart = args.q4_overlay_chart or (args.snapshot_dir / "visuals" / "q4_quarterly_overlay_indexed.png")
    add_q4_timing_overlay_slide(prs, tokens, q4_overlay_chart)

    # Methodology slide (Coverage Ratio calculation)
    add_methodology_slide(prs, tokens)
    
    # Save
    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.out)
    print(f"Deck saved to {args.out}")
    
    # Convert to PDF
    convert_to_pdf(args.out)

def convert_to_pdf(pptx_path):
    print(f"Converting {pptx_path} to PDF...")
    try:
        # soffice --headless --convert-to pdf <file> --outdir <dir>
        cmd = [
            "soffice", 
            "--headless", 
            "--convert-to", "pdf", 
            str(pptx_path), 
            "--outdir", str(pptx_path.parent)
        ]
        subprocess.run(cmd, check=True, capture_output=True)
        pdf_path = pptx_path.with_suffix(".pdf")
        print(f"PDF saved to {pdf_path}")
    except subprocess.CalledProcessError as e:
        print(f"Error converting to PDF: {e}")
    except FileNotFoundError:
        print("Error: 'soffice' command not found. Please install LibreOffice to enable PDF export.")

if __name__ == "__main__":
    main()
