"""
build_deck.py — CLI to assemble a branded, fully editable .pptx from CSVs & PNGs.

Inputs
------
- prd_answers.csv (or delta_summary.csv) with cohort-level metrics.
- result_sheets/*.png (event-study charts, histogram, bubble, etc.).
- *.potx brand template.

Usage
-----
python build_deck.py \
  --template /path/to/Premier-FY25-PPT-16x9.potx \
  --csv /path/to/prd_answers.csv \
  --charts /path/to/result_sheets \
  --out /path/to/output/DeckName.pptx \
  [--style-yaml /path/to/style_overrides.yaml]

The CLI auto-picks case studies that pass reconciliation gates (gap ≤ 30%, overlap ≥ 0.30),
and populates Headline, Method, Portfolio, Case Studies, Integrity, and Appendix headers.
"""
from dataclasses import replace
from pathlib import Path
from typing import Optional, List, Dict, Any, Iterable
import argparse
import json
import math
import re
import subprocess
from datetime import datetime
import shutil
import zipfile
import pandas as pd

from pptx import Presentation  # type: ignore
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER  # type: ignore

from src.selector import PortfolioStats, build_deck_variables, compute_portfolio_stats

from .style_tokens import load_tokens
from .charts_fallback import render_histogram, render_bubble
from .helpers import load_theme_colors
from .slides import (
    PortfolioStats,
    CaseStudy,
    create_title_slide,
    create_headline_findings,
    create_method_slide,
    create_portfolio_view,
    create_case_study,
    create_integrity_slide,
    create_section_header,
    ensure_contrast,
)

NUMERIC_SPLIT_RE = re.compile(r"(\d+)\.\s*\n\s*(\d)")

def _read_style_yaml(path: Optional[Path]) -> Optional[Dict[str, Any]]:
    if not path:
        return None
    try:
        import yaml  # type: ignore
    except Exception:
        raise RuntimeError("PyYAML is required to use --style-yaml")
    with open(path, "r") as f:
        return yaml.safe_load(f)

def parse_event_month(s: str) -> Optional[datetime]:
    try:
        return datetime.strptime(s, "%Y-%m")
    except Exception:
        return None


def _load_deck_vars_json(path: Path) -> Dict[str, Any]:
    with path.open("r", encoding="utf-8") as fh:
        return json.load(fh)


def _coerce_float(value: Any) -> Optional[float]:
    if value is None:
        return None
    try:
        return float(value)
    except (TypeError, ValueError):
        return None


def _coerce_int(value: Any) -> Optional[int]:
    if value is None:
        return None
    try:
        return int(float(value))
    except (TypeError, ValueError):
        return None


def _percent_or_fraction(value: Any) -> Optional[float]:
    val = _coerce_float(value)
    if val is None:
        return None
    if abs(val) > 1.0:
        return val / 100.0
    return val


def _stats_from_headline(headline: Dict[str, Any]) -> PortfolioStats:
    pct_sustained = _coerce_float(headline.get("pct_sustained"))
    portfolio_mean = _coerce_float(headline.get("portfolio_mean_pp"))
    recent_mean = _coerce_float(headline.get("recent_mean_pp"))
    direction_positive_pct = _coerce_float(headline.get("direction_positive_pct"))
    direction_positive_mean = _coerce_float(headline.get("direction_positive_mean_pp"))
    direction_negative_mean = _coerce_float(headline.get("direction_negative_mean_pp"))
    median_ttt = _coerce_float(headline.get("median_time_to_target_mo"))
    return PortfolioStats(
        pct_sustained=float(pct_sustained) / 100.0 if pct_sustained is not None else 0.0,
        median_time_to_target=float(median_ttt) if median_ttt is not None else None,
        portfolio_mean_delta=float(portfolio_mean) / 100.0 if portfolio_mean is not None else 0.0,
        eligible_count=int(_coerce_int(headline.get("eligible_total")) or 0),
        total_count=int(_coerce_int(headline.get("cohorts_total")) or 0),
        recent_mean_delta=float(recent_mean) / 100.0 if recent_mean is not None else None,
        direction_positive_pct=float(direction_positive_pct) / 100.0 if direction_positive_pct is not None else None,
        direction_positive_mean=float(direction_positive_mean) / 100.0 if direction_positive_mean is not None else None,
        direction_negative_mean=float(direction_negative_mean) / 100.0 if direction_negative_mean is not None else None,
    )


def _case_from_dict(payload: Dict[str, Any], case_type: str) -> Optional[CaseStudy]:
    program = str(payload.get("program_id") or "").strip()
    category = str(payload.get("category_id") or "").strip()
    if not program or not category:
        return None
    if payload.get("delta_pp") is not None:
        delta_raw = _coerce_float(payload.get("delta_pp")) or 0.0
        delta_fraction = delta_raw / 100.0
    else:
        delta_fraction = _percent_or_fraction(payload.get("delta_7_12")) or 0.0
    t0 = str(payload.get("t0") or payload.get("t0_event_month") or "").strip()
    members = _coerce_int(payload.get("members") or payload.get("N_members"))
    controls = _coerce_int(payload.get("controls") or payload.get("N_controls"))
    spend = _coerce_float(payload.get("member_t0_spend_6m"))
    gap = _coerce_float(payload.get("spend_gap_pct"))
    overlap = _coerce_float(payload.get("overlap"))
    dollars = _coerce_float(payload.get("dollars_shifted"))
    chart_path = payload.get("chart_path")
    integrity_flags = payload.get("integrity_flags", {}) if isinstance(payload.get("integrity_flags"), dict) else {}
    case = CaseStudy(
        program=program,
        category=category,
        delta_7_12=float(delta_fraction),
        t0_event_month=t0 or "",
        n_members=members if members is not None else 0,
        n_controls=controls if controls is not None else 0,
        spend_6m_usd=spend,
        recon_spend_gap_pct=gap,
        overlap_ratio=overlap,
        dollars_shifted=dollars,
        case_type=case_type,
        pretrend_risk_flag=bool(integrity_flags.get("pretrend_risk_flag")),
        placebo_flag=bool(integrity_flags.get("placebo_flag")),
        retention_ok=bool(integrity_flags.get("retention_ok", True)),
    )
    if chart_path:
        case.chart_path = str(chart_path)
    return case


def _cases_from_examples(examples: Dict[str, Any]) -> List[CaseStudy]:
    cases: List[CaseStudy] = []

    def _extend(items: Iterable[Dict[str, Any]], label: str, limit: Optional[int] = None) -> None:
        if not items:
            return
        selected = list(items)
        if limit is not None:
            selected = selected[:limit]
        for item in selected:
            case = _case_from_dict(item, label)
            if case:
                cases.append(case)

    _extend(examples.get("wins", []), "win", limit=2)
    neutral = examples.get("no_gain") or []
    if neutral:
        case = _case_from_dict(neutral[0], "no_gain")
        if case:
            cases.append(case)
    _extend(examples.get("losses", []), "loss")
    return cases


def find_chart(charts_dir: Path, program: str, category: str) -> Optional[Path]:
    if not charts_dir.exists():
        return None
    key = f"{program}_{category}".lower().replace(" ", "_")
    for p in charts_dir.glob("*.png"):
        name = p.stem.lower().replace(" ", "_")
        if key in name or all(k in name for k in [program.lower(), category.lower().replace(" ", "_")]):
            return p
    return None


def _remove_slide(prs: Presentation, index: int) -> None:
    slides = prs.slides
    slide_ids = slides._sldIdLst  # type: ignore[attr-defined]
    slide_id = slide_ids[index]
    rId = slide_id.rId
    prs.part.drop_rel(rId)
    slide_ids.remove(slide_id)


def _is_slide_empty(slide) -> bool:
    has_content = False
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return False
        if getattr(shape, "has_chart", False) or getattr(shape, "has_table", False):
            return False
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text or ""
            if text.strip():
                return False
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PLACEHOLDER:
            continue
        if getattr(shape, "shapes", None):
            return False
    return True


def _clear_initial_blanks(prs: Presentation) -> int:
    removed = 0
    while len(prs.slides) and _is_slide_empty(prs.slides[0]):
        _remove_slide(prs, 0)
        removed += 1
    return removed


def _prune_empty_slides(prs: Presentation) -> int:
    removed = 0
    for idx in reversed(range(len(prs.slides))):
        if _is_slide_empty(prs.slides[idx]):
            _remove_slide(prs, idx)
            removed += 1
    return removed


def _get_slide_title(slide) -> str:
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None:
        text = getattr(title_shape, "text", "")
        if text and text.strip():
            return text.strip()
    for shape in slide.shapes:
        text = getattr(shape, "text", "")
        if text and text.strip():
            return text.strip()
    return ""


def _is_empty_placeholder(shape) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    try:
        placeholder = shape.placeholder_format  # type: ignore[attr-defined]
    except ValueError:
        return False
    if placeholder.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.SUBTITLE):
        return False
    text = getattr(shape, "text", "") or ""
    return text.strip() == ""


def _remove_empty_placeholders(slide) -> int:
    removed = 0
    for idx in reversed(range(len(slide.shapes))):
        shape = slide.shapes[idx]
        if _is_empty_placeholder(shape):
            el = shape._element  # type: ignore[attr-defined]
            el.getparent().remove(el)
            removed += 1
    return removed


def _first_run_hex(shape) -> Optional[str]:
    if shape is None or not getattr(shape, "text_frame", None):
        return None
    tf = shape.text_frame
    if not tf.paragraphs:
        return None
    runs = tf.paragraphs[0].runs
    if not runs:
        return None
    run = runs[0]
    color = run.font.color.rgb
    if color is None:
        return None
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def _fix_numeric_wraps(slide) -> int:
    replacements = 0
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        tf = shape.text_frame
        for para in tf.paragraphs:
            if not para.runs:
                continue
            text = "".join(run.text for run in para.runs)
            new_text = NUMERIC_SPLIT_RE.sub(r"\1.\2", text)
            if new_text != text:
                replacements += 1
                first_run = para.runs[0]
                first_run.text = new_text
                for run in para.runs[1:]:
                    run.text = ""
    return replacements


def _count_empty_placeholders(prs: Presentation) -> int:
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            try:
                placeholder = shape.placeholder_format
            except ValueError:
                continue
            if placeholder.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.SUBTITLE):
                continue
            if not (shape.text or "").strip():
                count += 1
    return count


def _count_freeform_textboxes(slide) -> int:
    count = 0
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        try:
            shape.placeholder_format
            is_placeholder = True
        except ValueError:
            is_placeholder = False
        if not is_placeholder:
            count += 1
    return count


def _shapes_intersect_rect(a, b) -> bool:
    return not (a.left + a.width <= b.left or b.left + b.width <= a.left or a.top + a.height <= b.top or b.top + b.height <= a.top)


def _compute_badge_metrics(prs: Presentation) -> Dict[str, int]:
    overlaps = 0
    multiline = 0
    for idx in range(5, min(len(prs.slides), 10)):
        slide = prs.slides[idx]
        badges = [sh for sh in slide.shapes if sh.name == "delta_badge"]
        targets = []
        for sh in slide.shapes:
            if not getattr(sh, "has_text_frame", False):
                continue
            try:
                placeholder = sh.placeholder_format
            except ValueError:
                placeholder = None
            if placeholder is not None and sh.name != "delta_badge":
                targets.append(sh)
        for badge in badges:
            for target in targets:
                if _shapes_intersect_rect(badge, target):
                    overlaps += 1
        for badge in badges:
            if getattr(badge, "text_frame", None) and "\n" in badge.text_frame.text:
                multiline += 1
                break
    return {"badge_overlaps": overlaps, "badge_multiline": multiline}

def _materialize_template(template: Path, out_path: Path) -> Path:
    """
    Copy `template` to `out_path`, coercing POTX content-types to PPTX so python-pptx can load it.
    Returns the path that should be opened with Presentation().
    """
    if template.suffix.lower() == ".potx":
        with zipfile.ZipFile(template, "r") as src, zipfile.ZipFile(out_path, "w") as dst:
            for item in src.infolist():
                data = src.read(item.filename)
                if item.filename == "[Content_Types].xml":
                    data = data.replace(
                        b"application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
                        b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
                    )
                zi = zipfile.ZipInfo(filename=item.filename)
                zi.compress_type = item.compress_type
                zi.flag_bits = item.flag_bits
                zi.create_system = item.create_system
                zi.external_attr = item.external_attr
                zi.date_time = item.date_time
                dst.writestr(zi, data)
        return out_path
    else:
        if template != out_path:
            shutil.copyfile(template, out_path)
        return out_path


def build(args):
    csv_path = Path(args.csv)
    charts_dir = Path(args.charts)
    template = Path(args.template)
    out_path = Path(args.out)
    overrides = _read_style_yaml(Path(args.style_yaml)) if args.style_yaml else None

    df = pd.read_csv(csv_path)
    tokens = load_tokens(template, overrides)
    theme_cfg = {}
    if isinstance(overrides, dict):
        theme_cfg = overrides.get("theme", {}) or {}
    theme_colors = {}
    if theme_cfg.get("source"):
        theme_path = theme_cfg["source"]
        theme_colors = load_theme_colors(theme_path)
        print("theme_loaded=True", "theme_keys=" + str(sorted(theme_colors.keys())))
    else:
        print("theme_loaded=False")
    if theme_colors:
        tokens = replace(tokens, theme=theme_colors)

    target_fraction = (
        tokens.target_lift_pp / 100.0 if tokens.target_lift_pp and tokens.target_lift_pp > 1 else tokens.target_lift_pp
    )
    deck_vars_path = Path(args.deck_vars).expanduser().resolve() if getattr(args, "deck_vars", None) else None
    deck_vars_data: Optional[Dict[str, Any]] = None
    computed_deck_vars: Optional[Dict[str, Any]] = None
    if deck_vars_path:
        if deck_vars_path.exists():
            deck_vars_data = _load_deck_vars_json(deck_vars_path)
        else:
            print({"deck_vars_missing": str(deck_vars_path)})
    if deck_vars_data is None:
        computed_deck_vars = build_deck_variables(df, target_lift=target_fraction)
        deck_vars_data = computed_deck_vars
    headline_data = deck_vars_data.get("headline", {})
    target_override = _coerce_float(headline_data.get("target_lift_pp"))
    if target_override is not None:
        tokens = replace(tokens, target_lift_pp=float(target_override))
        target_fraction = (
            tokens.target_lift_pp / 100.0 if tokens.target_lift_pp and tokens.target_lift_pp > 1 else tokens.target_lift_pp
        )
    else:
        target_fraction = (
            tokens.target_lift_pp / 100.0 if tokens.target_lift_pp and tokens.target_lift_pp > 1 else tokens.target_lift_pp
        )
    if not deck_vars_data.get("examples"):
        if computed_deck_vars is None:
            computed_deck_vars = build_deck_variables(df, target_lift=target_fraction)
        deck_vars_data.setdefault("examples", computed_deck_vars.get("examples", {}))
        if not deck_vars_data.get("headline"):
            deck_vars_data["headline"] = computed_deck_vars.get("headline", {})
        headline_data = deck_vars_data.get("headline", {})

    out_path.parent.mkdir(parents=True, exist_ok=True)
    base_path = _materialize_template(template, out_path)
    prs = Presentation(str(base_path))
    initial_pruned = _clear_initial_blanks(prs)
    chosen_layouts: Dict[str, str] = {}
    removed_placeholders = 0
    wrap_fixes = 0
    contrast_fixed: List[int] = []
    contrast_colors: Dict[int, Optional[str]] = {}

    charts_dir.mkdir(parents=True, exist_ok=True)
    histogram_png = charts_dir / "portfolio_histogram.png"
    bubble_png = charts_dir / "portfolio_bubble.png"
    if not histogram_png.exists():
        render_histogram(df, histogram_png)
    if not bubble_png.exists():
        render_bubble(df, bubble_png)

    # Title
    slide = create_title_slide(prs, tokens, title="Performance-Group Awards: Market-Share Impact", subtitle=datetime.now().strftime("%B %d, %Y"), layout_choices=chosen_layouts)
    removed_placeholders += _remove_empty_placeholders(slide)
    wrap_fixes += _fix_numeric_wraps(slide)

    target_fraction = tokens.target_lift_pp / 100.0 if tokens.target_lift_pp and tokens.target_lift_pp > 1 else tokens.target_lift_pp

    # Headline
    stats = _stats_from_headline(headline_data)
    if (stats.total_count == 0 or stats.eligible_count == 0) and not df.empty:
        stats = compute_portfolio_stats(df, target_fraction)
    slide = create_headline_findings(prs, tokens, stats, layout_choices=chosen_layouts)
    removed_placeholders += _remove_empty_placeholders(slide)
    wrap_fixes += _fix_numeric_wraps(slide)

    # Method
    slide = create_method_slide(
        prs,
        tokens,
        target_lift_pp=headline_data.get("target_lift_pp"),
        recon_policy=deck_vars_data.get("recon_policy"),
        layout_choices=chosen_layouts,
    )
    removed_placeholders += _remove_empty_placeholders(slide)
    wrap_fixes += _fix_numeric_wraps(slide)

    # Portfolio view (optional inserts if images provided)
    histogram = histogram_png
    bubble = bubble_png
    slide = create_portfolio_view(
        prs,
        tokens,
        histogram_png=str(histogram) if histogram.exists() else None,
        bubble_png=str(bubble) if bubble.exists() else None,
        headline=headline_data,
        program_split=deck_vars_data.get("program_split"),
        drivers=deck_vars_data.get("drivers"),
        portfolio_stats=stats,
        layout_choices=chosen_layouts,
    )
    removed_placeholders += _remove_empty_placeholders(slide)
    wrap_fixes += _fix_numeric_wraps(slide)

    # Case studies
    slide = create_section_header(prs, tokens, "Case Studies", layout_choices=chosen_layouts)
    removed_placeholders += _remove_empty_placeholders(slide)
    wrap_fixes += _fix_numeric_wraps(slide)
    case_studies = _cases_from_examples(deck_vars_data.get("examples", {}))
    if not case_studies and computed_deck_vars is not None:
        case_studies = _cases_from_examples(computed_deck_vars.get("examples", {}))
    if not case_studies:
        fallback_cases = build_deck_variables(df, target_lift=target_fraction)
        case_studies = _cases_from_examples(fallback_cases.get("examples", {}))
    rendered_cases: List[CaseStudy] = []
    case_badges: List[Dict[str, Any]] = []
    case_geometries: List[Dict[str, Any]] = []
    for cs in case_studies:
        chart: Optional[Path] = None
        if cs.chart_path:
            candidate = Path(cs.chart_path)
            if not candidate.is_absolute():
                candidate_name = candidate.name
                candidate_from_dir = charts_dir / candidate_name
                candidate = candidate_from_dir if candidate_from_dir.exists() else charts_dir / candidate
            if candidate.exists():
                chart = candidate
        if chart is None:
            chart = find_chart(charts_dir, cs.program, cs.category)
        if not chart:
            print({"skip_case": {"program": cs.program, "category": cs.category}})
            continue
        cs.chart_path = str(chart)
        slide, badge_label, badge_positive = create_case_study(prs, tokens, cs, layout_choices=chosen_layouts)
        removed_placeholders += _remove_empty_placeholders(slide)
        wrap_fixes += _fix_numeric_wraps(slide)
        rendered_cases.append(cs)
        case_badges.append({
            "program": cs.program,
            "category": cs.category,
            "case_type": cs.case_type,
            "delta_pp": round(cs.delta_7_12 * 100.0, 1),
            "positive": bool(badge_positive),
            "label": badge_label,
            "warnings": {
                "pretrend_risk": cs.pretrend_risk_flag,
                "placebo": cs.placebo_flag,
                "retention_ok": cs.retention_ok,
            },
        })
        geometry = getattr(slide, "case_geometry", None)
        if geometry:
            case_geometries.append({
                "program": cs.program,
                "category": cs.category,
                "case_type": cs.case_type,
                "geometry": geometry,
            })

    # Integrity
    slide = create_integrity_slide(prs, tokens, layout_choices=chosen_layouts)
    removed_placeholders += _remove_empty_placeholders(slide)
    wrap_fixes += _fix_numeric_wraps(slide)

    # Appendix headers (partner packs are added in a subsequent pass, if desired)
    slide = create_section_header(prs, tokens, "Appendix — Manufacturer Chart Packs", layout_choices=chosen_layouts)
    removed_placeholders += _remove_empty_placeholders(slide)
    wrap_fixes += _fix_numeric_wraps(slide)

    for slide in prs.slides:
        removed_placeholders += _remove_empty_placeholders(slide)
        wrap_fixes += _fix_numeric_wraps(slide)

    pruned_after_build = _prune_empty_slides(prs)
    total_pruned = initial_pruned + pruned_after_build
    final_titles = [_get_slide_title(slide) for slide in prs.slides]
    title_shape = prs.slides[0].shapes.title if prs.slides else None
    headline_shape = prs.slides[1].shapes.title if len(prs.slides) > 1 else None
    title_rgb = _first_run_hex(title_shape)
    headline_rgb = _first_run_hex(headline_shape)
    method_metrics = getattr(prs.slides[2], "method_metrics", None) if len(prs.slides) > 2 else None
    portfolio_metrics = getattr(prs.slides[3], "portfolio_metrics", None) if len(prs.slides) > 3 else None
    headline_title_pt = getattr(prs.slides[1], "headline_title_pt", None) if len(prs.slides) > 1 else None
    headline_numbers_pt = getattr(prs.slides[1], "headline_numbers_pt", None) if len(prs.slides) > 1 else None
    headline_readability_ok = getattr(prs.slides[1], "headline_readability_ok", None) if len(prs.slides) > 1 else None
    empty_placeholder_count = _count_empty_placeholders(prs)
    title_freeforms = _count_freeform_textboxes(prs.slides[0]) if prs.slides else 0
    headline_freeforms = _count_freeform_textboxes(prs.slides[1]) if len(prs.slides) > 1 else 0
    pruned_slides_delta = pruned_after_build
    badge_metrics = _compute_badge_metrics(prs)
    summary = {
        "pct_sustained": stats.pct_sustained,
        "median_ttt": stats.median_time_to_target,
        "portfolio_mean_delta": stats.portfolio_mean_delta,
        "n_cases": len(rendered_cases),
        "target_lift_pp": tokens.target_lift_pp,
        "deck_vars_source": str(deck_vars_path) if deck_vars_path and deck_vars_path.exists() else ("computed" if getattr(args, "deck_vars", None) else None),
        "chosen_layouts": chosen_layouts,
        "slide_titles": final_titles,
        "cases_rendered": case_badges,
        "pruned_slides": total_pruned,
        "removed_empty_placeholders": removed_placeholders,
        "wrap_fixes": wrap_fixes,
        "title_rgb": title_rgb,
        "headline_rgb": headline_rgb,
        "headline_title_pt": headline_title_pt,
        "headline_numbers_pt": headline_numbers_pt,
        "headline_readability_ok": headline_readability_ok,
        "contrast_fixed_slides": contrast_fixed,
        "contrast_text_rgb": contrast_colors,
        "contrast_sample_rgb": next((color for color in contrast_colors.values() if color), None),
        "headline_layout": chosen_layouts.get("TITLE_CONTENT"),
        "empty_placeholders": empty_placeholder_count,
        "title_freeforms": title_freeforms,
        "headline_freeforms": headline_freeforms,
        "pruned_slides_delta": pruned_slides_delta,
    }
    if method_metrics:
        summary["how_we_measured"] = method_metrics
    if portfolio_metrics:
        summary["portfolio_view"] = portfolio_metrics
    if case_geometries:
        summary["case_studies_geometry"] = case_geometries
    summary.update(badge_metrics)
    summary["pdf_exported"] = None
    prs.save(str(out_path))
    pdf_path = _export_pdf(out_path)
    if pdf_path:
        summary["pdf_exported"] = pdf_path.name
    prop_summary = {
        "hl": summary.get("headline_layout"),
        "cfs": contrast_fixed,
        "csr": summary.get("contrast_sample_rgb"),
        "bo": summary.get("badge_overlaps"),
        "bm": summary.get("badge_multiline"),
        "ep": summary.get("empty_placeholders"),
        "tf": summary.get("title_freeforms"),
        "hf": summary.get("headline_freeforms"),
        "psd": summary.get("pruned_slides_delta"),
        "ht": summary.get("headline_title_pt"),
        "hn": summary.get("headline_numbers_pt"),
        "hr": summary.get("headline_readability_ok"),
        "pdf": summary.get("pdf_exported"),
    }
    # Persist summary metadata on the presentation for downstream audits.
    setattr(prs, "deck_summary", summary)
    setattr(prs, "headline_layout", summary.get("headline_layout"))
    try:
        import json
        prs.core_properties.comments = json.dumps(prop_summary, separators=(",", ":"))
    except Exception as exc:  # pragma: no cover - defensive persistence
        print({"deck_summary_write_failed": str(exc)})
    prs.save(str(out_path))
    print({"pruned_slides": total_pruned})
    print({"removed_empty_placeholders": removed_placeholders})
    print({"title_rgb": title_rgb, "headline_rgb": headline_rgb})
    print({"slide_titles": final_titles})
    print(summary)
    return out_path


def _export_pdf(pptx_path: Path) -> Optional[Path]:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        print({"pdf_export_skipped": "soffice not found"})
        return None
    out_dir = pptx_path.parent
    cmd = [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), str(pptx_path)]
    try:
        completed = subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    except subprocess.CalledProcessError as exc:
        stderr = (exc.stderr or b"").decode(errors="ignore").strip()
        print({"pdf_export_failed": stderr or str(exc)})
        return None
    pdf_path = out_dir / f"{pptx_path.stem}.pdf"
    if not pdf_path.exists():
        stdout = (completed.stdout or b"").decode(errors="ignore").strip()
        print({"pdf_export_missing": pdf_path.name, "stdout": stdout})
        return None
    return pdf_path

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--template", required=True, help=".potx brand template path")
    ap.add_argument("--csv", required=True, help="Path to prd_answers.csv or delta_summary.csv")
    ap.add_argument("--charts", required=True, help="Directory of chart PNGs")
    ap.add_argument("--out", required=True, help="Output .pptx path")
    ap.add_argument("--style-yaml", required=False, help="Optional style overrides YAML")
    ap.add_argument("--deck-vars", required=False, help="Optional path to deck_vars.json produced by the selector step")
    args = ap.parse_args()
    out = build(args)
    print(f"Wrote deck: {out}")

if __name__ == "__main__":
    main()
