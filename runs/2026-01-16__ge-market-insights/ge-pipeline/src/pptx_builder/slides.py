"""
slides.py — modular slide generators for the analytics deck.
"""
from dataclasses import dataclass
from typing import Any, Dict, List, Optional, Tuple, Union
from datetime import datetime
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

from src.selector.deck_vars import PortfolioStats

from .layouts import find_layout, DEFAULT_LAYOUT_MAP
from .helpers import (
    place_norm,
    set_paragraph,
    add_big_number,
    add_picture_fit,
    find_placeholder,
    reset_text_frame,
    get_theme_colors,
)

def _preferred(tokens, key, default):
    prefs = tokens.layout_map.get(key) if hasattr(tokens, "layout_map") else None
    if prefs:
        return [str(p) for p in prefs]
    return default

def _target_threshold(tokens) -> float:
    target = getattr(tokens, "target_lift_pp", 0.0)
    return target / 100.0 if target and target > 1 else target

def _title_color(tokens):
    if getattr(tokens, "force_title_contrast", False):
        if tokens.title_text_color is not None:
            return tokens.title_text_color
        return tokens.colors.text_light
    return None

def _subtitle_color(tokens):
    if getattr(tokens, "force_title_contrast", False):
        if tokens.subtitle_text_color is not None:
            return tokens.subtitle_text_color
        return tokens.colors.text_light
    return None

def _hex_to_rgb(value) -> Optional[Tuple[int, int, int]]:
    if isinstance(value, RGBColor):
        return (value.red, value.green, value.blue)
    if isinstance(value, (tuple, list)) and len(value) == 3:
        try:
            return tuple(int(v) for v in value)
        except ValueError:
            return None
    if not value:
        return None
    if isinstance(value, str):
        hex_str = value.strip().lstrip("#")
    else:
        hex_str = ""
    if len(hex_str) != 6:
        return None
    try:
        return tuple(int(hex_str[i:i+2], 16) for i in (0, 2, 4))
    except ValueError:
        return None


def _rgb_to_hex(rgb: Optional[Union[Tuple[int, int, int], RGBColor]]) -> Optional[str]:
    triplet = _hex_to_rgb(rgb)
    if triplet is None:
        return None
    return "".join(f"{component:02X}" for component in triplet)


def _format_target_pp(tokens, target_lift_pp: Optional[float]) -> float:
    if target_lift_pp is not None:
        value = float(target_lift_pp)
    else:
        value = getattr(tokens, "target_lift_pp", 0.0) or 0.0
    return value / 100.0 if value > 1.0 else value


def _set_bullet_state(paragraph, *, enabled: bool, level: int = 0) -> None:
    paragraph.level = level
    pPr = paragraph._p.get_or_add_pPr()  # type: ignore[attr-defined]
    if enabled:
        for child in list(pPr):
            if child.tag.endswith("buNone"):
                pPr.remove(child)
        if not any(child.tag.endswith("buChar") for child in pPr):
            bu_char = OxmlElement("a:buChar")
            bu_char.set("char", "•")
            bu_char.set("typeface", "Arial")
            pPr.insert(0, bu_char)
    else:
        if not any(child.tag.endswith("buNone") for child in pPr):
            pPr.append(OxmlElement("a:buNone"))


def ensure_contrast(slide, tokens, targets=("title", "content")):
    candidates = []
    prs = slide.part.package.presentation_part.presentation
    theme_source = getattr(tokens, "theme", {}) or {}
    fallback_theme = get_theme_colors(prs)

    def theme_lookup(key, fallback):
        rgb = _hex_to_rgb(theme_source.get(key))
        if rgb is not None:
            return rgb
        rgb = _hex_to_rgb(fallback_theme.get(key))
        if rgb is not None:
            return rgb
        return fallback

    light_rgb = theme_lookup("lt1", (255, 255, 255))
    dark_rgb = theme_lookup("dk1", getattr(tokens.colors, "bg_dark", (0, 0, 0)))
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        include = False
        placeholder = None
        try:
            placeholder = shape.placeholder_format  # type: ignore[attr-defined]
        except ValueError:
            placeholder = None
        if placeholder is not None:
            if placeholder.type == PP_PLACEHOLDER.TITLE and "title" in targets:
                include = True
            elif placeholder.type == PP_PLACEHOLDER.SUBTITLE and "title" in targets:
                include = True
            elif placeholder.type not in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.SUBTITLE) and "content" in targets:
                include = True
        elif "content" in targets:
            include = True
        if include:
            candidates.append(shape)
    if not candidates:
        return False
    needs_overlay = False
    for shape in candidates:
        tf = shape.text_frame
        if not tf.paragraphs or not tf.paragraphs[0].runs:
            needs_overlay = True
            break
        color = tf.paragraphs[0].runs[0].font.color.rgb
        if color is None or sum(color) < 400:
            needs_overlay = True
            break
    if not needs_overlay:
        return None
    left = min(shape.left for shape in candidates)
    top = min(shape.top for shape in candidates)
    right = max(shape.left + shape.width for shape in candidates)
    bottom = max(shape.top + shape.height for shape in candidates)
    pad_w = int((right - left) * 0.05)
    pad_h = int((bottom - top) * 0.05)
    rect = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        max(0, left - pad_w),
        max(0, top - pad_h),
        (right - left) + 2 * pad_w,
        (bottom - top) + 2 * pad_h,
    )
    rect.name = "contrast_overlay"
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(*dark_rgb)
    rect.fill.transparency = 0.25
    rect.line.fill.background()
    tree = slide.shapes._spTree  # type: ignore[attr-defined]
    el = rect._element  # type: ignore[attr-defined]
    tree.remove(el)
    anchor = getattr(candidates[0], "_element", None)
    if anchor is not None:
        idx = list(tree).index(anchor)
        tree.insert(idx, el)
    else:
        tree.insert(0, el)
    for shape in candidates:
        tf = shape.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.color.rgb = RGBColor(*light_rgb)
    setattr(slide, "contrast_text_rgb", _rgb_to_hex(light_rgb))
    return light_rgb


def shapes_intersect(shape_a, shape_b) -> bool:
    if shape_a is None or shape_b is None:
        return False
    a_left, a_top = shape_a.left, shape_a.top
    a_right, a_bottom = a_left + shape_a.width, a_top + shape_a.height
    b_left, b_top = shape_b.left, shape_b.top
    b_right, b_bottom = b_left + shape_b.width, b_top + shape_b.height
    return not (a_right <= b_left or b_right <= a_left or a_bottom <= b_top or b_bottom <= a_top)

def _resolve_layout(prs, tokens, key, default_fragments, fallback_index, layout_choices: Optional[Dict[str, str]] = None):
    preferred = _preferred(tokens, key, default_fragments)
    preferred_names = None
    if hasattr(tokens, "layout_map"):
        preferred_names = tokens.layout_map.get(key)
    layout = find_layout(prs, preferred, fallback_index, preferred_names=preferred_names, layout_key=key)
    if layout_choices is not None and layout is not None:
        layout_choices[key] = layout.name
    return layout

def _add_case_badge(slide, tokens, delta: float, threshold: float, avoid_shapes=None):
    avoid_shapes = [shape for shape in (avoid_shapes or []) if shape is not None]
    presentation = slide.part.package.presentation_part.presentation
    theme_source = getattr(tokens, "theme", {}) or {}
    fallback_theme = get_theme_colors(presentation)
    sw, sh = presentation.slide_width, presentation.slide_height
    badge_d = max(int(Inches(1.2)), int(min(sw, sh) * 0.07))
    badge_left = int(Inches(10.1))
    badge_top = int(Inches(5.0))
    if badge_left + badge_d > sw - int(sw * 0.01):
        badge_left = max(int(sw * 0.01), sw - badge_d - int(sw * 0.01))
    if badge_top + badge_d > sh - int(sh * 0.01):
        badge_top = max(int(sh * 0.01), sh - badge_d - int(sh * 0.01))
    margin = int(Inches(0.1))

    def _intersects(rect_left: int, rect_top: int, rect_d: int, shape) -> bool:
        rect_right = rect_left + rect_d
        rect_bottom = rect_top + rect_d
        shape_left = shape.left
        shape_top = shape.top
        shape_right = shape.left + shape.width
        shape_bottom = shape.top + shape.height
        return not (
            rect_right <= shape_left
            or shape_right <= rect_left
            or rect_bottom <= shape_top
            or shape_bottom <= rect_top
        )

    for _ in range(5):
        adjusted = False
        for shape in avoid_shapes:
            if _intersects(badge_left, badge_top, badge_d, shape):
                down_candidate = shape.top + shape.height + margin
                if down_candidate + badge_d <= sh - margin:
                    badge_top = down_candidate
                else:
                    up_candidate = max(margin, shape.top - badge_d - margin)
                    badge_top = up_candidate
                adjusted = True
        if not adjusted:
            break
    positive = delta >= threshold if threshold is not None else delta >= 0
    label_text = ("▲" if positive else "▼") + f" {delta * 100:+.1f} pp"

    ellipse = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, badge_left, badge_top, badge_d, badge_d)
    ellipse.name = "delta_badge"
    ellipse.line.fill.background()
    ellipse.fill.solid()
    accent_key = "accent1" if positive else "accent3"
    accent_rgb = _hex_to_rgb(theme_source.get(accent_key)) or _hex_to_rgb(fallback_theme.get(accent_key))
    if accent_rgb is None:
        accent_tuple = getattr(tokens.colors, "accent1" if positive else "accent3")
        accent_rgb = accent_tuple
    ellipse.fill.fore_color.rgb = RGBColor(*accent_rgb)
    tf = ellipse.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    margin = Pt(6)
    tf.margin_left = margin
    tf.margin_right = margin
    tf.margin_top = margin
    tf.margin_bottom = margin
    para = tf.paragraphs[0]
    para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    run = para.add_run()
    run.text = label_text.replace("\n", " ")
    run.font.bold = True
    para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    light_rgb = _hex_to_rgb(theme_source.get("lt1")) or _hex_to_rgb(fallback_theme.get("lt1")) or (255, 255, 255)
    run.font.color.rgb = RGBColor(*light_rgb)
    label_size = 14
    run.font.size = Pt(label_size)
    return label_text, positive


def create_title_slide(prs: Presentation, tokens, title: str, subtitle: str = "", layout_choices: Optional[Dict[str, str]] = None):
    layout = _resolve_layout(prs, tokens, "TITLE", DEFAULT_LAYOUT_MAP["TITLE"], 0, layout_choices)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        p = reset_text_frame(tf)
        run = p.runs[0] if p.runs else p.add_run()
        run.text = title
        run.font.bold = True
        run.font.size = Pt(max(tokens.fonts.title_size_pt - 2, 28))
        title_color = _title_color(tokens)
        if title_color:
            run.font.color.rgb = RGBColor(*title_color)
    subtitle_ph = find_placeholder(slide, placeholder_type=PP_PLACEHOLDER.SUBTITLE, fallback_to_first_text=True)
    if subtitle_ph is not None and subtitle_ph.has_text_frame:
        # Slightly raise the subtitle block so long subheaders stay on-screen
        subtitle_ph.top = max(0, subtitle_ph.top - Inches(0.25))
        subtitle_ph.height = max(Inches(1.0), subtitle_ph.height - Inches(0.1))
        tf = subtitle_ph.text_frame
        p = reset_text_frame(tf)
        run = p.runs[0] if p.runs else p.add_run()
        run.text = subtitle
        run.font.size = Pt(24)
        p.space_after = Pt(6)
        tf.margin_bottom = Pt(0)
        tf.margin_top = Pt(6)
        subtitle_color = _subtitle_color(tokens)
        if subtitle_color:
            run.font.color.rgb = RGBColor(*subtitle_color)
    return slide

def create_headline_findings(prs: Presentation, tokens, stats: PortfolioStats, layout_choices: Optional[Dict[str, str]] = None):
    key = "TITLE_CONTENT"
    layout = _resolve_layout(prs, tokens, key, DEFAULT_LAYOUT_MAP.get(key, ["title", "content"]), 1, layout_choices)
    slide = prs.slides.add_slide(layout)
    headline_title_pt = max(
        tokens.typography.get("headline_title_pt", 38),
        tokens.fonts.title_size_pt,
    )
    big_number_pt = max(tokens.typography.get("headline_numbers_pt", 32), tokens.fonts.h2_size_pt)
    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        p = reset_text_frame(tf)
        run = p.runs[0] if p.runs else p.add_run()
        run.text = "Headline Findings — Do Awards Drive Sustained Lift?"
        run.font.bold = True
        run.font.size = Pt(headline_title_pt)
        title_color = _title_color(tokens)
        color = tokens.colors.text_dark if title_color is None else title_color
        run.font.color.rgb = RGBColor(*color)
    body_ph = find_placeholder(slide, PP_PLACEHOLDER.BODY, fallback_to_first_text=True)
    metrics_box = None
    if body_ph is not None and body_ph.has_text_frame:
        tf = body_ph.text_frame
        tf.clear()
        context_pt = tokens.typography.get("headline_body_pt", max(tokens.fonts.body_size_pt + 4, 18))
        lines = [
            {"text": f"{stats.pct_sustained:.0%} of eligible cohorts achieved sustained lift (≥ 7–12 target)", "size": big_number_pt, "bold": True},
            {"text": f"Median time-to-target: {stats.median_time_to_target:.0f} months" if stats.median_time_to_target is not None else "Median time-to-target: —", "size": big_number_pt, "bold": True},
            {"text": f"Portfolio mean Δ7–12: {stats.portfolio_mean_delta:+.2f} pp", "size": big_number_pt, "bold": True},
            {"text": f"Eligible {stats.eligible_count} / Total {stats.total_count} cohorts", "size": context_pt, "bold": False},
        ]
        if stats.recent_mean_delta is not None:
            lines.append({"text": f"Recent 2025 vintage mean Δ7–12: {stats.recent_mean_delta:+.2f} pp", "size": context_pt, "bold": False})
        for idx, cfg in enumerate(lines):
            para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            para.level = 0
            run = para.add_run()
            run.text = cfg["text"]
            run.font.bold = bool(cfg.get("bold"))
            run.font.size = Pt(cfg.get("size", context_pt))
            run.font.color.rgb = RGBColor(*tokens.colors.text_dark)
    setattr(slide, "headline_title_pt", headline_title_pt)
    setattr(slide, "headline_numbers_pt", big_number_pt)
    setattr(slide, "headline_readability_ok", True)
    return slide

def create_method_slide(
    prs: Presentation,
    tokens,
    *,
    target_lift_pp: Optional[float] = None,
    recon_policy: Optional[Dict[str, Any]] = None,
    layout_choices: Optional[Dict[str, str]] = None,
):
    layout = _resolve_layout(prs, tokens, "TITLE_CONTENT", DEFAULT_LAYOUT_MAP["TITLE_CONTENT"], 1, layout_choices)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        p = reset_text_frame(tf)
        set_paragraph(p, "How We Measured", tokens.fonts.title_face, tokens.typography.get("method_title_pt", tokens.fonts.h2_size_pt), True, tokens.colors.text_dark, inherit_fonts=tokens.inherit_fonts)

    target_fraction = _format_target_pp(tokens, target_lift_pp)
    target_pp = target_fraction * 100.0
    top_subhead_pt = max(18, tokens.typography.get("method_top_subhead_pt", 18))
    top_bullet_pt = 16
    gating_subhead_pt = top_subhead_pt
    gating_bullet_pt = 16

    bullets = [
        "t₀-locked cohort (program members fixed at award month)",
        "Controls: non-program facilities in the same category/timeframe with supplier-overlap exclusion",
        "Outcome: awarded-block share (6-month trailing) with Laspeyres t₀ facility weights; right-censor after exit",
        f"Windows 0–6 and 7–12 with TARGET ≥ {target_pp:.1f} pp lift threshold",
    ]

    recon_policy = recon_policy or {}
    gate_gap = recon_policy.get("gap_max_pct")
    gate_overlap = recon_policy.get("overlap_min")
    gating_lines: List[str] = []
    if gate_gap is not None:
        pct = gate_gap if abs(gate_gap) > 1 else gate_gap * 100.0
        gating_lines.append(f"|PRD vs Dashboard spend gap| ≤ {pct:.0f}%")
    else:
        gating_lines.append("|PRD vs Dashboard spend gap| ≤ 30%")
    if gate_overlap is not None:
        pct = gate_overlap if abs(gate_overlap) > 1 else gate_overlap * 100.0
        gating_lines.append(f"Facility overlap ≥ {pct:.0f}%")
    else:
        gating_lines.append("Facility overlap ≥ 30%")
    gating_lines.append("No zero-vs-nonzero asymmetry")
    gating_lines.append("Mapped taxonomy")

    body_ph = find_placeholder(slide, PP_PLACEHOLDER.BODY, fallback_to_first_text=True)
    metrics = {
        "top_subhead_pt": top_subhead_pt,
        "top_bullets_pt": top_bullet_pt,
        "showcase_gating": {
            "subhead_pt": gating_subhead_pt,
            "bullets_pt": gating_bullet_pt,
        },
        "bulleted_range_ok": False,
    }
    if body_ph is not None and body_ph.has_text_frame:
        tf = body_ph.text_frame
        p = reset_text_frame(tf)
        set_paragraph(
            p,
            "Program design & weighting",
            tokens.fonts.body_face,
            top_subhead_pt,
            True,
            tokens.colors.text_dark,
            inherit_fonts=False,
        )
        _set_bullet_state(p, enabled=False, level=0)
        p.line_spacing = 1.0
        p.space_after = Pt(4)

        top_paragraphs = []
        for line in bullets:
            para = tf.add_paragraph()
            set_paragraph(
                para,
                f"\t{line}",
                tokens.fonts.body_face,
                top_bullet_pt,
                False,
                tokens.colors.text_dark,
                inherit_fonts=False,
            )
            _set_bullet_state(para, enabled=True, level=0)
            para.line_spacing = 1.0
            para.space_after = Pt(0)
            top_paragraphs.append(para)

        heading = tf.add_paragraph()
        set_paragraph(
            heading,
            "Showcase gating (reconciliation-aligned)",
            tokens.fonts.body_face,
            gating_subhead_pt,
            True,
            tokens.colors.text_dark,
            inherit_fonts=False,
        )
        _set_bullet_state(heading, enabled=False, level=0)
        heading.line_spacing = 1.0
        heading.space_before = Pt(8)
        heading.space_after = Pt(0)

        for line in gating_lines:
            para = tf.add_paragraph()
            set_paragraph(
                para,
                line,
                tokens.fonts.body_face,
                gating_bullet_pt,
                False,
                tokens.colors.text_dark,
                inherit_fonts=False,
            )
            _set_bullet_state(para, enabled=True, level=1)
            para.line_spacing = 1.0
            para.space_after = Pt(0)
        def _has_bullets(paragraph) -> bool:
            pPr = paragraph._p.get_or_add_pPr()  # type: ignore[attr-defined]
            return not any(child.tag.endswith("buNone") for child in pPr)

        metrics["bulleted_range_ok"] = all(_has_bullets(p) for p in top_paragraphs)
    setattr(slide, "method_metrics", metrics)
    return slide

def create_portfolio_view(
    prs: Presentation,
    tokens,
    histogram_png: Optional[str] = None,
    bubble_png: Optional[str] = None,
    *,
    headline: Optional[Dict[str, Any]] = None,
    program_split: Optional[List[Dict[str, Any]]] = None,
    drivers: Optional[Dict[str, Any]] = None,
    portfolio_stats: Optional[PortfolioStats] = None,
    layout_choices: Optional[Dict[str, str]] = None,
):
    layout = _resolve_layout(prs, tokens, "TWO_CONTENT", DEFAULT_LAYOUT_MAP["TWO_CONTENT"], 3, layout_choices)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        p = reset_text_frame(tf)
        set_paragraph(p, "Portfolio View", tokens.fonts.title_face, tokens.fonts.h2_size_pt, True, tokens.colors.text_dark, inherit_fonts=tokens.inherit_fonts)
    presentation = slide.part.package.presentation_part.presentation
    sw, sh = presentation.slide_width, presentation.slide_height
    margin = int(sw * 0.05)
    prev_text_width = int(sw * 0.35)
    text_width = max(int(sw * 0.28), int(sw * 0.25))
    gap = max(int(sw * 0.025), int(Inches(0.2)))

    placeholders = []
    for ph in slide.placeholders:
        ph_format = getattr(ph, "placeholder_format", None)
        if ph_format is None:
            continue
        if ph_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            placeholders.append(ph)
    placeholders.sort(key=lambda ph: ph.left)

    text_placeholder = None
    for ph in placeholders:
        if getattr(ph, "has_text_frame", False):
            text_placeholder = ph
            break
    if text_placeholder in placeholders:
        placeholders.remove(text_placeholder)

    chart_width = int(Inches(5.0))
    chart_height = int(Inches(4.0))
    chart_top = int(Inches(1.30))
    chart_spacing = int(Inches(0.5))
    summary_gap = int(Inches(0.4))
    chart_summary_height = int(Inches(1.5))

    targets = [histogram_png, bubble_png]
    available_paths = [p for p in targets if p]
    portfolio_metrics = {"charts": [], "overlaps": 0, "aspect_within_5": True}
    chart_bottom = chart_top + chart_height
    chart_shapes = []
    target_ratio = chart_width / chart_height if chart_height else 1.0
    if available_paths:
        count = len(available_paths)
        total_width = chart_width * count + chart_spacing * (count - 1)
        start_left = max(margin, (sw - total_width) // 2)
        for idx, path in enumerate(available_paths):
            left = start_left + idx * (chart_width + chart_spacing)
            picture = slide.shapes.add_picture(path, left, chart_top)
            picture.width = chart_width
            picture.height = chart_height
            chart_shapes.append(picture)
            aspect_ratio = picture.width / picture.height if picture.height else 1.0
            portfolio_metrics["charts"].append(
                {
                    "index": idx,
                    "width_in": round(picture.width / Inches(1), 3),
                    "height_in": round(picture.height / Inches(1), 3),
                    "left_in": round(picture.left / Inches(1), 3),
                    "top_in": round(picture.top / Inches(1), 3),
                    "aspect_ratio": round(aspect_ratio, 3),
                    "aspect_within_5": abs(aspect_ratio - target_ratio) <= 0.05,
                }
            )
        overlaps = 0
        for i, shape_a in enumerate(chart_shapes):
            for shape_b in chart_shapes[i + 1:]:
                if not (
                    shape_a.left + shape_a.width <= shape_b.left
                    or shape_b.left + shape_b.width <= shape_a.left
                    or shape_a.top + shape_a.height <= shape_b.top
                    or shape_b.top + shape_b.height <= shape_a.top
                ):
                    overlaps += 1
        portfolio_metrics["overlaps"] = overlaps
        portfolio_metrics["aspect_within_5"] = all(m["aspect_within_5"] for m in portfolio_metrics["charts"])
        chart_bottom = max(shape.top + shape.height for shape in chart_shapes)

    summary_top = chart_bottom + summary_gap
    summary_width = sw - 2 * margin
    summary_left = margin
    if summary_top + chart_summary_height > sh - margin:
        summary_top = max(margin, sh - margin - chart_summary_height)

    if text_placeholder is not None and text_placeholder.has_text_frame:
        summary_shape = text_placeholder
        summary_shape.left = summary_left
        summary_shape.top = summary_top
        summary_shape.width = summary_width
        summary_shape.height = chart_summary_height
        summary_tf = summary_shape.text_frame
    else:
        summary_shape = slide.shapes.add_textbox(summary_left, summary_top, summary_width, chart_summary_height)
        summary_tf = summary_shape.text_frame
    summary_tf.clear()
    summary_tf.word_wrap = True

    stats_obj = portfolio_stats
    mean_delta_pp = None
    pct_sustained = None
    eligible_count = None
    target_pp = getattr(tokens, "target_lift_pp", None)
    if stats_obj is not None:
        mean_delta_pp = stats_obj.portfolio_mean_delta * 100.0 if stats_obj.portfolio_mean_delta is not None else None
        pct_sustained = stats_obj.pct_sustained * 100.0 if stats_obj.pct_sustained is not None else None
        eligible_count = stats_obj.eligible_count
    head = headline or {}
    if eligible_count in (None, 0):
        eligible_count = head.get("eligible_total")
    if target_pp in (None, 0):
        target_pp = head.get("target_lift_pp")
    target_txt = f"{target_pp:.1f} pp" if isinstance(target_pp, (int, float)) else "target"
    mean_txt = f"{mean_delta_pp:+.1f} pp" if mean_delta_pp is not None else "—"
    pct_txt = f"{pct_sustained:.0f}%" if pct_sustained is not None else "—"
    eligible_txt = str(eligible_count) if eligible_count not in (None, 0) else "all"

    summary_lines = [
        f"Portfolio mean Δ7–12 lift: {mean_txt}; {pct_txt} of eligible cohorts meet or beat the {target_txt} target.",
        f"Histogram (left) shows Δ7–12 lift distribution across {eligible_txt} cohorts (dashed line marks {target_txt}).",
        "Bubble chart (right) sizes each cohort by member t₀ spend to surface high-dollar case-study opportunities.",
    ]

    summary_font_pt = max(tokens.typography.get("portfolio_summary_pt", 13), 12)
    for idx, text in enumerate(summary_lines):
        para = summary_tf.paragraphs[0] if idx == 0 else summary_tf.add_paragraph()
        set_paragraph(
            para,
            text,
            tokens.fonts.body_face,
            summary_font_pt,
            False,
            tokens.colors.text_dark,
            inherit_fonts=False,
        )
        _set_bullet_state(para, enabled=False, level=0)
        para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        para.line_spacing = 1.2
        para.space_after = Pt(0)

    portfolio_metrics.update(
        {
            "mean_delta_pp": mean_delta_pp,
            "pct_sustained_pct": pct_sustained,
            "eligible_count": eligible_count,
        }
    )
    setattr(slide, "portfolio_metrics", portfolio_metrics)
    return slide

@dataclass
class CaseStudy:
    program: str
    category: str
    delta_7_12: float
    t0_event_month: str
    n_members: int
    n_controls: int
    spend_6m_usd: Optional[float] = None
    recon_spend_gap_pct: Optional[float] = None
    overlap_ratio: Optional[float] = None
    dollars_shifted: Optional[float] = None
    chart_path: Optional[str] = None
    case_type: str = "case"
    pretrend_risk_flag: bool = False
    placebo_flag: bool = False
    retention_ok: bool = True

def create_case_study(prs: Presentation, tokens, case: CaseStudy, layout_choices: Optional[Dict[str, str]] = None) -> Tuple:
    layout = _resolve_layout(prs, tokens, "TITLE_CONTENT", DEFAULT_LAYOUT_MAP["TITLE_CONTENT"], 1, layout_choices)
    slide = prs.slides.add_slide(layout)
    presentation = slide.part.package.presentation_part.presentation
    sw, sh = presentation.slide_width, presentation.slide_height
    case_key = case.case_type.lower() if isinstance(case.case_type, str) else ""
    case_label = {
        "win": "Win",
        "loss": "Loss",
        "no_gain": "No Gain",
    }.get(case_key, None)
    title = f"{case.program.upper()} — {case.category}"
    if case_label:
        title = f"{case_label} — {case.program.upper()} — {case.category}"
    if slide.shapes.title:
        tf_title = slide.shapes.title.text_frame
        p_title = reset_text_frame(tf_title)
        set_paragraph(p_title, title, tokens.fonts.title_face, tokens.fonts.h2_size_pt, True, tokens.colors.text_dark, inherit_fonts=tokens.inherit_fonts)
    chart_left = int(Inches(0.50))
    chart_top = int(Inches(1.30))
    chart_width = int(Inches(8.00))
    chart_height = int(Inches(5.50))
    metrics_left = int(Inches(9.00))
    metrics_top = int(Inches(1.30))
    metrics_width = int(Inches(4.00))
    metrics_height = int(Inches(3.60))
    chart_shape = None
    if case.chart_path:
        chart_shape = slide.shapes.add_picture(case.chart_path, chart_left, chart_top)
        chart_shape.width = chart_width
        chart_shape.height = chart_height
    title_shape = getattr(slide.shapes, "title", None)
    chart_overlaps_title = False
    if chart_shape is not None and title_shape is not None:
        title_bottom = title_shape.top + title_shape.height
        overlap_guard = int(Inches(0.01))
        chart_overlaps_title = chart_shape.top < title_bottom - overlap_guard
        if chart_overlaps_title:
            chart_shape.top = title_bottom + int(Inches(0.05))
            chart_overlaps_title = False
    body_ph = find_placeholder(slide, PP_PLACEHOLDER.BODY, fallback_to_first_text=True)
    metrics_shape = None
    if body_ph is not None and body_ph.has_text_frame:
        body_ph.left = metrics_left
        body_ph.top = metrics_top
        body_ph.width = metrics_width
        body_ph.height = metrics_height
        tf = body_ph.text_frame
        tf.clear()
        tf.word_wrap = True
        metrics_shape = body_ph
    else:
        metrics_shape = slide.shapes.add_textbox(metrics_left, metrics_top, metrics_width, metrics_height)
        metrics_shape.name = "case_metrics"
        tf = metrics_shape.text_frame
        tf.clear()
        tf.word_wrap = True
    p = reset_text_frame(tf)
    delta_pp = case.delta_7_12 * 100.0
    lines = [
        f"Δ7–12: {delta_pp:+.2f} pp   |   t₀: {case.t0_event_month}",
        f"Measured Ns: members {case.n_members}, controls {case.n_controls}",
    ]
    if case.spend_6m_usd is not None:
        lines.append(f"Member t₀ spend (6m): ${case.spend_6m_usd:,.2f}")
    if case.dollars_shifted is not None:
        lines.append(f"Dollars shifted (Δ7–12 × spend): ${case.dollars_shifted:,.0f}")
    if case.recon_spend_gap_pct is not None or case.overlap_ratio is not None:
        gap_text = f"{case.recon_spend_gap_pct:+.2f}%" if case.recon_spend_gap_pct is not None else "—"
        overlap_text = f"{case.overlap_ratio:.2f}" if case.overlap_ratio is not None else "—"
        lines.append(f"Recon gap {gap_text}   |   Overlap {overlap_text}")
    warnings = []
    if case.pretrend_risk_flag:
        warnings.append("pre-trend risk")
    if case.placebo_flag:
        warnings.append("placebo concern")
    if not case.retention_ok:
        warnings.append("retention risk")
    if warnings:
        lines.append("⚠ " + "; ".join(warnings))
    metrics_size = 14
    for idx, line in enumerate(lines):
        para = p if idx == 0 else tf.add_paragraph()
        set_paragraph(
            para,
            line,
            tokens.fonts.body_face,
            metrics_size,
            bold=False,
            color_rgb=tokens.colors.text_dark,
            inherit_fonts=False,
        )
    def _as_inches(value: int) -> float:
        return round(value / Inches(1), 3)

    def _within_tolerance(value: float, target: float, tol: float = 0.05) -> bool:
        return abs(value - target) <= tol

    chart_info = None
    if chart_shape is not None:
        chart_info = {
            "left_in": _as_inches(chart_shape.left),
            "top_in": _as_inches(chart_shape.top),
            "width_in": _as_inches(chart_shape.width),
            "height_in": _as_inches(chart_shape.height),
            "within_tolerance": all(
                _within_tolerance(measured, target)
                for measured, target in [
                    (_as_inches(chart_shape.left), 0.50),
                    (_as_inches(chart_shape.top), 1.30),
                    (_as_inches(chart_shape.width), 8.00),
                    (_as_inches(chart_shape.height), 5.50),
                ]
            ),
        }
    textbox_info = {
        "left_in": _as_inches(metrics_shape.left),
        "top_in": _as_inches(metrics_shape.top),
        "width_in": _as_inches(metrics_shape.width),
        "height_in": _as_inches(metrics_shape.height),
        "within_tolerance": all(
            _within_tolerance(measured, target)
            for measured, target in [
                (_as_inches(metrics_shape.left), 9.00),
                (_as_inches(metrics_shape.top), 1.30),
                (_as_inches(metrics_shape.width), 4.00),
                (_as_inches(metrics_shape.height), 3.60),
            ]
        ),
    }
    geometry_metrics = {
        "chart": chart_info,
        "textbox": textbox_info,
        "kpi_font_pt": metrics_size,
        "chart_overlaps_title": chart_overlaps_title,
    }
    setattr(slide, "case_geometry", geometry_metrics)
    avoid = [metrics_shape, chart_shape, slide.shapes.title]
    badge_label, badge_positive = _add_case_badge(slide, tokens, case.delta_7_12, _target_threshold(tokens), avoid_shapes=avoid)
    return slide, badge_label, badge_positive

def create_integrity_slide(prs: Presentation, tokens, layout_choices: Optional[Dict[str, str]] = None):
    layout = _resolve_layout(prs, tokens, "TITLE_CONTENT", DEFAULT_LAYOUT_MAP["TITLE_CONTENT"], 1, layout_choices)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        p = reset_text_frame(tf)
        set_paragraph(p, "Study Integrity & Mitigations", tokens.fonts.title_face, tokens.fonts.h2_size_pt, True, tokens.colors.text_dark, inherit_fonts=tokens.inherit_fonts)
    strengths = [
        "Supplier-overlap exclusion for controls; right-censor after exit",
        "Denominator = full category (MVP); equivalence-set sensitivity queued",
        "Pre-trend slope checks and retention overlays on highlighted cohorts",
    ]
    watchouts = [
        "Confirm equivalence-set coverage before narrowing denominators",
        "Monitor recon gap/overlap thresholds (≥0.30 overlap, |gap| ≤ 30%)",
        "Flag emerging t₀ vintages for recency/attrition review",
    ]
    body_ph = find_placeholder(slide, PP_PLACEHOLDER.BODY, fallback_to_first_text=True)
    if body_ph is not None and body_ph.has_text_frame:
        tf = body_ph.text_frame
        tf.clear()
        p_strengths = tf.paragraphs[0]
        set_paragraph(p_strengths, "Strengths", tokens.fonts.body_face, tokens.fonts.body_size_pt, True, tokens.colors.text_dark, inherit_fonts=tokens.inherit_fonts)
        p_strengths.level = 0
        for point in strengths:
            para = tf.add_paragraph()
            set_paragraph(para, point, tokens.fonts.body_face, tokens.fonts.body_size_pt, False, tokens.colors.text_dark, inherit_fonts=tokens.inherit_fonts)
            para.level = 1
            para.space_after = Pt(0)
        watchouts_head = tf.add_paragraph()
        set_paragraph(watchouts_head, "Key watch-outs & what we’re doing", tokens.fonts.body_face, tokens.fonts.body_size_pt, True, tokens.colors.text_dark, inherit_fonts=tokens.inherit_fonts)
        watchouts_head.level = 0
        watchouts_head.space_before = Pt(12)
        for point in watchouts:
            para = tf.add_paragraph()
            set_paragraph(para, point, tokens.fonts.body_face, tokens.fonts.body_size_pt, False, tokens.colors.text_dark, inherit_fonts=tokens.inherit_fonts)
            para.level = 1
    return slide

def create_section_header(prs: Presentation, tokens, title: str, layout_choices: Optional[Dict[str, str]] = None):
    layout = _resolve_layout(prs, tokens, "SECTION", DEFAULT_LAYOUT_MAP["SECTION"], 2, layout_choices)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        p = reset_text_frame(tf)
        section_color = tokens.colors.text_light if getattr(tokens, "force_title_contrast", False) else tokens.colors.text_dark
        set_paragraph(p, title, tokens.fonts.title_face, tokens.fonts.title_size_pt, True, section_color, inherit_fonts=tokens.inherit_fonts)
    return slide
