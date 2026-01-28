"""Deck QA helper checks for presentation validation."""
from __future__ import annotations

import json
from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation  # type: ignore
from pptx.util import Inches


def _portfolio_chart_sizes(prs: Presentation) -> List[float]:
    if len(prs.slides) < 4:
        return []
    slide = prs.slides[3]
    sizes = []
    for shape in slide.shapes:
        if shape.shape_type == 13:  # picture
            sizes.append(round(shape.width / Inches(1), 3))
    return sizes


def run(deck_path: Path | str) -> Dict[str, Any]:
    deck_path = Path(deck_path)
    prs = Presentation(str(deck_path))
    width, height = prs.slide_width, prs.slide_height

    blank_titles: List[int] = []
    offslide_shapes: List[Dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides):
        title = slide.shapes.title
        if title is not None:
            text = title.text_frame.text if title.has_text_frame else ""
            if not (text or "").strip():
                blank_titles.append(idx + 1)
        for shape in slide.shapes:
            if not (getattr(shape, "has_text_frame", False) or shape.shape_type in {13, 1}):
                continue
            left = getattr(shape, "left", 0)
            top = getattr(shape, "top", 0)
            right = left + getattr(shape, "width", 0)
            bottom = top + getattr(shape, "height", 0)
            if left < 0 or top < 0 or right > width or bottom > height:
                offslide_shapes.append(
                    {
                        "slide": idx + 1,
                        "name": getattr(shape, "name", ""),
                        "bounds_in": (
                            round(left / Inches(1), 3),
                            round(top / Inches(1), 3),
                            round(right / Inches(1), 3),
                            round(bottom / Inches(1), 3),
                        ),
                    }
                )

    comments = prs.core_properties.comments
    prop_summary: Dict[str, Any] = {}
    if comments:
        try:
            prop_summary = json.loads(comments)
        except json.JSONDecodeError:
            pass

    results: Dict[str, Any] = {
        "blank_titles": blank_titles,
        "offslide_shapes": offslide_shapes,
        "badge_overlaps": prop_summary.get("bo"),
        "badge_multiline": prop_summary.get("bm"),
        "contrast_fixed_slides": prop_summary.get("cfs"),
        "contrast_sample_rgb": prop_summary.get("csr"),
        "portfolio_chart_widths_in": _portfolio_chart_sizes(prs),
    }
    return results


def main() -> None:
    import argparse

    parser = argparse.ArgumentParser(description="Deck QA checks")
    parser.add_argument("deck", type=Path, help="Path to deck pptx")
    args = parser.parse_args()
    output = run(args.deck)
    print(output)


if __name__ == "__main__":
    main()
