import os
from pptx import Presentation

RUN_ID = "2026-01-06__accelerating_technology_delivery_presentation"
DRAFT_FILE = f"runs/{RUN_ID}/exports/draft_v4_content.pptx"

if not os.path.exists(DRAFT_FILE):
    print("Draft file not found.")
    exit(1)

prs = Presentation(DRAFT_FILE)
print(f"Inspecting {DRAFT_FILE} ({len(prs.slides)} slides)")

for i, slide in enumerate(prs.slides):
    print(f"Slide {i} Layout: {slide.slide_layout.name}")
    for shape in slide.placeholders:
        print(f"  Placeholder [{shape.placeholder_format.idx}] '{shape.name}': HasTextFrame={shape.has_text_frame}")
        if shape.has_text_frame and shape.text:
            print(f"    TEXT: {shape.text[:50]}...")
            if shape.text_frame.paragraphs:
                p = shape.text_frame.paragraphs[0]
                if p.font.color.type:
                     print(f"    COLOR (Para): {p.font.color.rgb}")
                else:
                     print(f"    COLOR (Para): None (Inherited)")
        else:
            print(f"    (Empty or No TextFrame)")

