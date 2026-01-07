import os
from pptx import Presentation

# Configuration
RUN_ID = "2026-01-06__accelerating_technology_delivery_presentation"
INPUT_TEMPLATE = f"runs/{RUN_ID}/inputs/brand/Premier_Template.pptx"
OUTPUT_LIBRARY = f"runs/{RUN_ID}/assets/Premier_Components.pptx"

def create_library():
    print(f"Creating Component Library from {INPUT_TEMPLATE}...")
    if not os.path.exists(INPUT_TEMPLATE):
        raise FileNotFoundError(f"Template not found: {INPUT_TEMPLATE}")

    prs = Presentation(INPUT_TEMPLATE)
    
    # We want to KEEP only the layouts we need and SAVE them as a 'Library' file 
    # that generate_via_components.py can load and clone from.
    # Actually, simpler: We just instantiate the template, add sample slides, and save.
    # But generate_via_components needs to know WHICH layout index maps to WHICH component logic.
    
    # MAPPING (Based on manual inspection of Premier-FY25-PPT):
    # Layout 0:  Title Slide
    # Layout 1:  Title and Content (Standard)
    # Layout 4:  Section Header (Assumed)
    
    # Let's create a visual reference deck to confirm
    
    # --- COMPONENT 1: TITLE SLIDE (ID: COMP_TITLE) ---
    slide_layout = prs.slide_layouts[0] 
    slide = prs.slides.add_slide(slide_layout)
    # Mark it for reference
    if slide.shapes.title: slide.shapes.title.text = "{{TITLE}}"
    try:
        if len(slide.placeholders) > 1: slide.placeholders[1].text = "{{SUBTITLE}}"
    except: pass
    
    # --- COMPONENT 2: AGENDA / LIST (ID: COMP_LIST) ---
    slide_layout = prs.slide_layouts[1] # Standard Content
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title: slide.shapes.title.text = "{{SECTION_TITLE}}"
    try:
        if len(slide.placeholders) > 1: slide.placeholders[1].text = "{{BULLET_LIST}}"
    except: pass
    
    # --- COMPONENT 3: SECTION HEADER (ID: COMP_SECTION) ---
     # Usually 'Section Header' is further down. Let's guess index 8 for "Premier" template or finding one by name?
    # For fail-safety, we will use Layout 1 again if 8 doesn't exist, but usually templates have ~11 layouts.
    layout_index = 8 if len(prs.slide_layouts) > 8 else 1
    slide_layout = prs.slide_layouts[layout_index] 
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title: slide.shapes.title.text = "{{SECTION_HEADER}}"
    
    os.makedirs(os.path.dirname(OUTPUT_LIBRARY), exist_ok=True)
    prs.save(OUTPUT_LIBRARY)
    print(f"Library saved to {OUTPUT_LIBRARY}")

if __name__ == "__main__":
    create_library()
