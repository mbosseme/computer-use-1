import sys
import os
import json
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from shapely.geometry import box
import rtree

def get_shape_bounds(shape):
    """
    Returns (left, bottom, right, top) for shapely box, 
    but PPT uses (left, top, width, height) where y increases downwards.
    Shapely uses Cartesian (y increases upwards).
    To map PPT to Shapely:
    x = left
    y = -top (flip y axis)
    """
    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height
    
    # Inverted Y axis for cartesian compatibility
    # Bottom in cartesian is -(top + height)
    # Top in cartesian is -top
    return (left, -(top + height), left + width, -top)

def get_normalized_bounds(shape, slide_width, slide_height):
    """
    Returns bounds in 0-1000 normalized space.
    Format: [left, top, width, height]
    """
    if slide_width == 0 or slide_height == 0:
        return [0, 0, 0, 0]
        
    left_norm = int((shape.left / slide_width) * 1000)
    top_norm = int((shape.top / slide_height) * 1000)
    width_norm = int((shape.width / slide_width) * 1000)
    height_norm = int((shape.height / slide_height) * 1000)
    
    return [left_norm, top_norm, width_norm, height_norm]

def lint_presentation(pptx_path, map_output_path=None):
    print(f"Linting {pptx_path}...")
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"Error loading presentation: {e}")
        return False
    
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    issues = []
    shape_map = {} # Mapping of slide_idx -> list of shapes
    
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        print(f"Checking Slide {slide_num}...")
        
        shape_map[i] = []
        
        shapes = list(slide.shapes)
        spatial_index = rtree.index.Index()
        shape_geoms = {}
        
        # Pass 1: Build Index & Map
        for idx, shape in enumerate(shapes):
            # Add to map regardless of rotation
            shape_info = {
                "id": shape.shape_id,
                "name": shape.name,
                "bbox": get_normalized_bounds(shape, slide_width, slide_height),
                "text": shape.text if hasattr(shape, "text") else ""
            }
            shape_map[i].append(shape_info)

            if shape.rotation != 0:
                # Rotation creates complex polygons, skipping for MVP
                continue
                
            bounds = get_shape_bounds(shape) # (minx, miny, maxx, maxy)
            geom = box(*bounds)
            shape_geoms[idx] = geom
            spatial_index.insert(idx, bounds)
            
            # Check 1: Off-canvas
            # Allow some tolerance?
            slide_box = box(0, -slide_height, slide_width, 0)
            if not slide_box.contains(geom):
                # Tolerance check (floating point of 100 EMUs which is tiny, maybe increase?)
                # 360000 EMU = 1 cm roughly
                if not slide_box.buffer(360000).contains(geom):
                    issues.append(f"Slide {slide_num}: Shape '{shape.name}' is partially off-slide.")

        # Pass 2: Check Overlaps
        for idx, shape in enumerate(shapes):
            if idx not in shape_geoms: continue
            
            geom = shape_geoms[idx]
            possible_overlaps = list(spatial_index.intersection(geom.bounds))
            
            for other_idx in possible_overlaps:
                if idx == other_idx: continue
                if other_idx not in shape_geoms: continue
                
                other_geom = shape_geoms[other_idx]
                
                if geom.intersects(other_geom):
                    intersection_area = geom.intersection(other_geom).area
                    min_area = min(geom.area, other_geom.area)
                    
                    # Ignore tiny overlaps (frames touching)
                    if min_area > 0 and (intersection_area / min_area) > 0.05:
                        # 5% overlap threshold
                        other_shape = shapes[other_idx]
                        
                        # Filter out innocuous overlaps?
                        # e.g. Title usually doesn't overlap Body, but sometimes frames touch.
                        # Pictures often overlap backgrounds.
                        
                        # Just reporting for now.
                        issues.append(f"Slide {slide_num}: Overlap detected between '{shape.name}' and '{other_shape.name}'")

    # Export map if requested
    if map_output_path:
        with open(map_output_path, "w") as f:
            json.dump(shape_map, f, indent=2)
        print(f"Shape map exported to {map_output_path}")

    if issues:
        print("\n--- LINT FAILURES DETECTED ---")
        for issue in issues:
            print(f"[FAIL] {issue}")
        print(f"Total Issues: {len(issues)}")
        return False
    else:
        print("\n[PASS] No geometry issues found.")
        return True

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 lint_slides.py <input.pptx> [--export-map <path>]")
        sys.exit(1)
    
    pptx_path = sys.argv[1]
    map_path = None
    
    if len(sys.argv) > 3 and sys.argv[2] == "--export-map":
        map_path = sys.argv[3]
        
    success = lint_presentation(pptx_path, map_path)
    sys.exit(0 if success else 1)
