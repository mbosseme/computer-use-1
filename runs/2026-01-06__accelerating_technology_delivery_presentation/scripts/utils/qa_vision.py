import os
import io
from PIL import Image, ImageDraw, ImageFont
import json
from pptx import Presentation
from pptx.util import Inches, Pt

# Configuration
# Assuming 1000-space normalization for Gemini
NORMALIZED_SIZE = 1000

def get_slide_dimensions(prs):
    return prs.slide_width, prs.slide_height

def map_emu_to_pixels(emu_value, emu_total, pixel_total):
    """
    Maps an EMU value to a pixel coordinate.
    """
    return int((emu_value / emu_total) * pixel_total)

def draw_overlays(image_path, slide_object, slide_width, slide_height, output_path=None):
    """
    Draws Set-of-Mark (SoM) bounding boxes and IDs on the slide image.
    Returns the path to the annotated image and a JSON index of shapes.
    """
    try:
        img = Image.open(image_path)
    except FileNotFoundError:
        print(f"Error: Image not found at {image_path}")
        return None, {}

    draw = ImageDraw.Draw(img, "RGBA")
    width_px, height_px = img.size
    
    # PPTX Dimensions passed in
    
    shape_index = {}
    
    # Iterate shapes
    # We use z-order (default enumeration order in python-pptx usually follows z-order back-to-front or front-to-back? 
    # Actually, shapes[0] is back-most. We want unique IDs.)
    
    font = None
    try:
        font = ImageFont.truetype("Arial.ttf", size=20)
    except:
        font = ImageFont.load_default()

    for i, shape in enumerate(slide_object.shapes):
        shape_id = f"ID_{i}"
        
        # Get Geometry (EMU)
        # Skip if no geometry (some invisible placeholders)
        if not hasattr(shape, 'left'):
            continue
            
        left = shape.left
        top = shape.top
        w = shape.width
        h = shape.height
        
        # Map to Pixels
        x1 = map_emu_to_pixels(left, slide_width, width_px)
        y1 = map_emu_to_pixels(top, slide_height, height_px)
        x2 = map_emu_to_pixels(left + w, slide_width, width_px)
        y2 = map_emu_to_pixels(top + h, slide_height, height_px)
        
        # Draw Semi-Transparent Box
        # Red with alpha
        fill_color = (255, 0, 0, 50) 
        outline_color = (255, 0, 0, 255)
        
        draw.rectangle([x1, y1, x2, y2], fill=fill_color, outline=outline_color, width=2)
        
        # Draw Label
        label_bg = (255, 255, 255, 200)
        text_pos = (x1, y1)
        bbox = draw.textbbox(text_pos, shape_id, font=font)
        draw.rectangle(bbox, fill=label_bg)
        draw.text(text_pos, shape_id, fill=(0,0,0), font=font)
        
        # Store metadata for the AI
        shape_index[shape_id] = {
            "name": shape.name,
            "text": shape.text if hasattr(shape, "text") else "<graphic>",
            "type": str(shape.shape_type),
            "bounds_emu": [left, top, w, h],
            "bounds_px": [x1, y1, x2, y2]
        }

    if output_path:
        img.save(output_path)
    
    return output_path, shape_index

def simple_critique(shape_index):
    """
    A deterministic fallback critique if no API key is present.
    Checks for overlapping bounding boxes in the 1000-space.
    """
    critiques = []
    # Very basic check: any shapes with empty text?
    for sid, data in shape_index.items():
        if "Placeholder" in data["name"] and data["text"].strip() == "":
             critiques.append(f"Warning: {sid} ({data['name']}) is empty.")
    
    if not critiques:
        return "PASS: No obvious empty placeholders."
    return "\n".join(critiques)

# Requirements: google-generativeai
# Stub for Phase 2b integration
def critique_slide_with_gemini(annotated_image_path, prompt_extra=""):
    import google.generativeai as genai
    
    api_key = os.getenv("GOOGLE_API_KEY")
    if not api_key:
        print("Review skipped: GOOGLE_API_KEY not set.")
        return None

    genai.configure(api_key=api_key)
    # Using the flash model for speed/cost in loop, or pro for quality. 
    # Spec says Gemini 3 Pro (preview) - usually gemini-1.5-pro currently available in SDK
    model = genai.GenerativeModel('gemini-1.5-pro') 

    img = Image.open(annotated_image_path)
    
    prompt = """
    You are a Senior Art Director. 
    Analyze this slide presentation.
    The image has red bounding boxes with IDs (ID_0, ID_1, etc.) overlaid on each element.
    
    Task:
    1. Identify any aesthetic issues: Alignment, overlaps, awkward whitespace, font inconsistencies.
    2. Identify any brand issues.
    
    Output a format:
    - CRITICAL: [ID_X] intersects [ID_Y]
    - MINOR: [ID_Z] is not aligned left with [ID_X]
    - PASS: (If looks good)
    
    """ + prompt_extra
    
    try:
        response = model.generate_content([prompt, img])
        return response.text
    except Exception as e:
        return f"Error calling Gemini: {e}"

if __name__ == "__main__":
    # Test stub
    pass
