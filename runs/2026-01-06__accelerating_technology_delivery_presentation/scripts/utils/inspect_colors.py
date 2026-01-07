from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE

def inspect_slide_1(pptx_path):
    print(f"Inspecting: {pptx_path}")
    prs = Presentation(pptx_path)
    slide = prs.slides[0] # Slide 1
    
    print(f"--- Slide 1 Inspection ---")
    for i, shape in enumerate(slide.shapes):
        print(f"\n[ID_{i}] Name: {shape.name}")
        if not shape.has_text_frame:
            print("  (No TextFrame)")
            continue
            
        print(f"  Text Content: '{shape.text}'")
        
        # Check Paragraph/Run properties
        if shape.text_frame.paragraphs:
            p = shape.text_frame.paragraphs[0]
            if p.runs:
                r = p.runs[0]
                print(f"  Font used: {r.font.name}")
                if r.font.color and r.font.color.type:
                    print(f"  Color Type: {r.font.color.type}")
                    if r.font.color.type == MSO_COLOR_TYPE.RGB:
                        print(f"  RGB: {r.font.color.rgb}")
                    elif r.font.color.type == MSO_COLOR_TYPE.SCHEME:
                         print(f"  Theme Color: {r.font.color.theme_color}")
                else:
                    print("  Color: Default/Inherited (likely Black if not styled)")

if __name__ == "__main__":
    inspect_slide_1("runs/2026-01-06__accelerating_technology_delivery_presentation/assets/Premier_Components.pptx")
