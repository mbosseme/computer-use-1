import os
from pptx import Presentation

RUN_ID = "2026-01-06__accelerating_technology_delivery_presentation"
COMPONENT_LIBRARY = f"runs/{RUN_ID}/assets/Premier_Components.pptx"

if not os.path.exists(COMPONENT_LIBRARY):
    print(f"Error: {COMPONENT_LIBRARY} not found.")
    exit(1)

prs = Presentation(COMPONENT_LIBRARY)

print(f"Loaded {COMPONENT_LIBRARY}")
print(f"Total Slides (Prototypes): {len(prs.slides)}")
print("-" * 40)

for i, slide in enumerate(prs.slides):
    layout = slide.slide_layout
    print(f"Slide {i} (Layout: {layout.name})")
    for ph in slide.placeholders:
        print(f"  Placeholder [{ph.placeholder_format.idx}]: '{ph.name}' (Type: {ph.placeholder_format.type})")
        if hasattr(ph, "text") and ph.text:
            print(f"    Current Text: {ph.text[:50]}...")
    print("-" * 40)
