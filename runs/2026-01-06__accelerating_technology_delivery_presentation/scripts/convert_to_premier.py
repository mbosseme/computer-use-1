from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt


RUN_ROOT = Path(__file__).resolve().parents[1]


def delete_slide(prs: Presentation, index: int) -> None:
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[index])


def remove_unused_placeholders(slide) -> None:
    to_remove = []
    for shape in slide.placeholders:
        if shape.has_text_frame:
            if not shape.text or not shape.text.strip():
                to_remove.append(shape)
    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)


def style_text_frame(
    text_frame,
    font_color: Optional[RGBColor] = None,
    margin_left=None,
    margin_top=None,
    *,
    font_size_pt: Optional[float] = None,
    bold: Optional[bool] = None,
    auto_size: Optional[MSO_AUTO_SIZE] = None,
    word_wrap: Optional[bool] = None,
) -> None:
    if margin_left is not None:
        text_frame.margin_left = margin_left
    if margin_top is not None:
        text_frame.margin_top = margin_top

    if word_wrap is not None:
        text_frame.word_wrap = word_wrap
    if auto_size is not None:
        text_frame.auto_size = auto_size

    for p in text_frame.paragraphs:
        if font_color:
            p.font.color.rgb = font_color
        if font_size_pt is not None:
            p.font.size = Pt(font_size_pt)
        if bold is not None:
            p.font.bold = bold

        for run in p.runs:
            if font_color:
                run.font.color.rgb = font_color
            if font_size_pt is not None:
                run.font.size = Pt(font_size_pt)
            if bold is not None:
                run.font.bold = bold


def ph(slide, idx: int):
    # Placeholder lookup by placeholder_format.idx (not list index)
    for p in slide.placeholders:
        try:
            if p.placeholder_format.idx == idx:
                return p
        except Exception:
            continue
    return None


def _bbox_norm_to_emu(bbox: List[int], slide_w: int, slide_h: int) -> Tuple[int, int, int, int]:
    x, y, w, h = _clamp_bbox_norm(bbox)
    left = int((x / 1000) * slide_w)
    top = int((y / 1000) * slide_h)
    width = int((w / 1000) * slide_w)
    height = int((h / 1000) * slide_h)
    return left, top, width, height


def _crop_bbox_px(img: Image.Image, bbox: List[int]) -> Image.Image:
    x, y, w, h = _clamp_bbox_norm(bbox)
    left = int((x / 1000) * img.size[0])
    top = int((y / 1000) * img.size[1])
    right = int(((x + w) / 1000) * img.size[0])
    bottom = int(((y + h) / 1000) * img.size[1])

    left = max(0, min(left, img.size[0] - 1))
    top = max(0, min(top, img.size[1] - 1))
    right = max(left + 1, min(right, img.size[0]))
    bottom = max(top + 1, min(bottom, img.size[1]))

    return img.crop((left, top, right, bottom))


def _clamp_bbox_norm(bbox: List[int]) -> List[int]:
    """Clamp bbox [x,y,w,h] to 0-1000 normalized space and ensure positive size."""
    if not bbox or len(bbox) != 4:
        return [0, 0, 1, 1]
    x, y, w, h = [int(v) for v in bbox]

    # Normalize negatives
    x = max(0, x)
    y = max(0, y)
    w = max(1, w)
    h = max(1, h)

    # Clamp right/bottom
    if x + w > 1000:
        w = max(1, 1000 - x)
    if y + h > 1000:
        h = max(1, 1000 - y)

    # Extra safety: ensure within bounds
    x = min(999, x)
    y = min(999, y)
    return [x, y, w, h]


def _inset_bbox_norm(bbox: List[int], inset: int = 6) -> List[int]:
    """Inset bbox on all sides (normalized units) to reduce overlaps."""
    x, y, w, h = _clamp_bbox_norm(bbox)
    if w <= inset * 2 + 1 or h <= inset * 2 + 1:
        return [x, y, w, h]
    return _clamp_bbox_norm([x + inset, y + inset, w - inset * 2, h - inset * 2])


def _bbox_area_ratio(bbox: List[int]) -> float:
    x, y, w, h = _clamp_bbox_norm(bbox)
    return (w * h) / 1_000_000.0


def _join_bullets(bullets: Any) -> str:
    if not bullets:
        return ""
    if isinstance(bullets, str):
        return bullets
    if isinstance(bullets, list):
        return "\n".join([str(b).strip() for b in bullets if str(b).strip()])
    return str(bullets)


def _layout_select(layout_hint: str, has_subtitle: bool, layouts: Dict[str, Any]):
    hint = (layout_hint or "").strip().lower()
    if hint == "title":
        return layouts["title"]
    if hint == "section":
        return layouts["section"]
    if hint == "two_col":
        return layouts["two_col_subtitle"] if has_subtitle else layouts["two_col"]
    # Default
    return layouts["content"]


def _bboxes_overlap(a: List[int], b: List[int]) -> bool:
    ax, ay, aw, ah = _clamp_bbox_norm(a)
    bx, by, bw, bh = _clamp_bbox_norm(b)
    return (ax < bx + bw) and (ax + aw > bx) and (ay < by + bh) and (ay + ah > by)


def _clamp_bbox_to_region(bbox: List[int], region: List[int]) -> List[int]:
    """Clamp bbox within region, both in normalized [0..1000] units.

    region: [x, y, w, h]
    """
    x, y, w, h = _clamp_bbox_norm(bbox)
    rx, ry, rw, rh = _clamp_bbox_norm(region)
    # Ensure bbox not larger than region
    w = min(w, rw)
    h = min(h, rh)
    x = max(rx, min(x, rx + rw - w))
    y = max(ry, min(y, ry + rh - h))
    return [x, y, w, h]


def _grid_layout_bboxes(count: int, region: List[int], pad: int = 16) -> List[List[int]]:
    """Lay out count boxes in a non-overlapping grid within region.

    This is intentionally simple and deterministic to satisfy geometry lint.
    """
    if count <= 0:
        return []

    rx, ry, rw, rh = _clamp_bbox_norm(region)

    # Choose a compact grid. Allow more columns to reduce rows (taller boxes),
    # which helps reduce text clipping for dense diagrams.
    cols = int(max(1, min(6, round(count ** 0.5))))
    while cols * ((count + cols - 1) // cols) < count:
        cols += 1
    rows = int((count + cols - 1) // cols)

    cell_w = max(1, rw // cols)
    cell_h = max(1, rh // rows)
    box_w = max(1, cell_w - pad * 2)
    box_h = max(1, cell_h - pad * 2)

    out: List[List[int]] = []
    for i in range(count):
        r = i // cols
        c = i % cols
        x = rx + c * cell_w + pad
        y = ry + r * cell_h + pad
        out.append(_clamp_bbox_to_region([x, y, box_w, box_h], region))
    return out


def _add_diagram_nodes(
    slide,
    spec: Dict[str, Any],
    slide_w: int,
    slide_h: int,
    white: RGBColor,
    *,
    min_top_emu: int = 0,
):
    diagram = spec.get("diagram") or {}
    nodes = diagram.get("nodes") or []
    if not isinstance(nodes, list):
        return

    candidate_labels: List[str] = []
    candidate_bboxes: List[List[int]] = []

    # Keep nodes but avoid duplicating obvious template text.
    for node in nodes:
        try:
            label = (node.get("label") or "").strip()
            label_l = label.lower()
            bbox = node.get("bbox")
            if not label or not bbox or len(bbox) != 4:
                continue
            if label_l.startswith("title") or label_l.startswith("subtitle") or "title block" in label_l:
                continue
            candidate_labels.append(label)
            candidate_bboxes.append(_inset_bbox_norm(bbox, inset=10))
        except Exception:
            continue

    if not candidate_bboxes:
        return

    # Preserve a rough reading order (top-to-bottom, left-to-right).
    items = list(zip(candidate_labels, candidate_bboxes))
    items.sort(key=lambda it: (_clamp_bbox_norm(it[1])[1], _clamp_bbox_norm(it[1])[0]))
    candidate_labels = [it[0] for it in items]
    candidate_bboxes = [it[1] for it in items]

    min_top_norm = int((max(0, min_top_emu) / slide_h) * 1000)
    # Keep diagram content well below the title region to avoid overlaps.
    top = max(320, min_top_norm)
    bottom_limit = 880  # keep content above footer/logo area
    safe_region = [90, top, 820, max(1, bottom_limit - top)]
    # Deterministic grid layout to avoid overlaps.
    resolved = _grid_layout_bboxes(len(candidate_bboxes), safe_region, pad=24)

    for label, bbox in zip(candidate_labels, resolved):
        left, top, width, height = _bbox_norm_to_emu(bbox, slide_w, slide_h)
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left,
            top,
            width,
            height,
        )
        if shape.has_text_frame:
            shape.text_frame.text = label
            style_text_frame(
                shape.text_frame,
                font_color=white,
                font_size_pt=9,
                auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
                word_wrap=True,
            )


def build_deck(
    component_library: Path,
    specs_dir: Path,
    source_renders_dir: Path,
    output_pptx: Path,
    crops_dir: Path,
) -> None:
    prs = Presentation(str(component_library))
    original_slide_count = len(prs.slides)

    # Layouts (based on known stable indices from the run)
    layouts = {
        "title": prs.slides[0].slide_layout,
        "section": prs.slides[1].slide_layout,
        "content": prs.slides[5].slide_layout,  # Dark 1a
        "two_col": prs.slides[13].slide_layout,  # Dark 1b
        "two_col_subtitle": prs.slides[9].slide_layout,  # Dark Subtitle 1a
        "blank": prs.slides[15].slide_layout,  # Premier Blank Dark 1a
    }

    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)

    WHITE = RGBColor(255, 255, 255)

    spec_paths = sorted(specs_dir.glob("slide_*.json"))
    if not spec_paths:
        raise SystemExit(f"No specs found in {specs_dir}")

    crops_dir.mkdir(parents=True, exist_ok=True)

    for spec_path in spec_paths:
        spec = json.loads(spec_path.read_text(encoding="utf-8"))
        slide_index = int((spec.get("_meta") or {}).get("slide_index") or int(spec_path.stem.split("_")[1]))

        title = (spec.get("title") or "").strip()
        subtitle = (spec.get("subtitle") or "").strip()
        bullets = _join_bullets(spec.get("bullets"))
        layout_hint = (spec.get("layout_hint") or "content").strip()

        photos = spec.get("photos") or []

        diagram_nodes = ((spec.get("diagram") or {}).get("nodes") or [])
        has_diagram = isinstance(diagram_nodes, list) and len(diagram_nodes) > 0

        # Only render diagram nodes on slides without photos.
        # Otherwise we tend to create container rectangles that overlap the cropped photos,
        # which the geometry linter flags.
        use_diagram_nodes = has_diagram and not bool(photos)

        # If a photo touches the top edge, many Premier layouts have a title placeholder
        # spanning into that region, which creates overlaps in Gate 1. Use a blank layout
        # and place text manually instead.
        def _photo_touches_top(ph_list: list[dict]) -> bool:
            for p in ph_list:
                bbox = p.get("bbox")
                if not bbox or len(bbox) != 4:
                    continue
                try:
                    _x, y, _w, _h = bbox
                except Exception:
                    continue
                if int(y) <= 20:
                    return True
            return False

        # For photo slides, prefer blank to avoid placeholder/photo intersections.
        needs_blank_for_photos = bool(photos)

        # Diagram slides: always start from blank to avoid placeholder collisions.
        # Photo-to-top slides: blank avoids title-placeholder overlap.
        base_layout = (
            layouts["blank"]
            if (use_diagram_nodes or needs_blank_for_photos)
            else _layout_select(layout_hint, bool(subtitle), layouts)
        )

        slide = prs.slides.add_slide(base_layout)

        # Title
        p0 = ph(slide, 0)
        title_shape_for_bottom = None
        if title:
            if p0 and not needs_blank_for_photos:
                p0.text = title
                style_text_frame(
                    p0.text_frame,
                    font_color=WHITE,
                    auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
                    word_wrap=True,
                )
                title_shape_for_bottom = p0
            else:
                # Blank layout: place title on left side to avoid photo collisions.
                title_top = Inches(0.8)
                title_height = Inches(1.1)
                if needs_blank_for_photos and photos:
                    # For card-style slides where photos start near the top, keep the title tighter.
                    try:
                        min_photo_y = min(int(p.get("bbox")[1]) for p in photos if p.get("bbox") and len(p.get("bbox")) == 4)
                    except Exception:
                        min_photo_y = 1000
                    if min_photo_y < 240:
                        title_top = Inches(0.35)
                        title_height = Inches(0.85)
                title_box = slide.shapes.add_textbox(
                    Inches(0.7),
                    title_top,
                    Inches(5.6),
                    title_height,
                )
                title_box.text_frame.text = title
                style_text_frame(
                    title_box.text_frame,
                    font_color=WHITE,
                    auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
                    word_wrap=True,
                )
                title_shape_for_bottom = title_box

        # Subtitle (where applicable; our known layouts use idx 10)
        p10 = ph(slide, 10)
        if subtitle and (not use_diagram_nodes):
            if p10 and not needs_blank_for_photos:
                # In some layouts, idx 10 is body. We'll only place subtitle here if slide_hint asks for subtitle.
                if layout_hint.strip().lower() in {"title", "two_col"}:
                    p10.text = subtitle
                    style_text_frame(
                        p10.text_frame,
                        font_color=WHITE,
                        margin_left=Inches(0.5),
                        auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
                        word_wrap=True,
                    )
            else:
                subtitle_box = slide.shapes.add_textbox(
                    Inches(0.75),
                    Inches(2.0),
                    Inches(5.4),
                    Inches(0.6),
                )
                subtitle_box.text_frame.text = subtitle
                style_text_frame(
                    subtitle_box.text_frame,
                    font_color=WHITE,
                    auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
                    word_wrap=True,
                )

        # Body text (only if the layout has an obvious body placeholder and we're not doing a diagram reconstruction)
        if (not use_diagram_nodes) and p10 and bullets and (not needs_blank_for_photos) and layout_hint.strip().lower() not in {"title"}:
            # If we already used p10 for subtitle (two_col_subtitle layout also has subtitle), try to keep body in columns.
            if layout_hint.strip().lower() == "two_col":
                left = ph(slide, 1)
                if left:
                    left.text = bullets
                    style_text_frame(
                        left.text_frame,
                        font_color=WHITE,
                        margin_top=Inches(0.15),
                        auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
                        word_wrap=True,
                    )
            else:
                # Default content body
                p10.text = bullets
                style_text_frame(
                    p10.text_frame,
                    font_color=WHITE,
                    margin_top=Inches(0.15),
                    auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
                    word_wrap=True,
                )

        # Blank-layout slides with photos: add body text manually so we avoid placeholder collisions.
        # If a photo occupies the upper-left region (e.g., card comparisons), place the body below it.
        if (not use_diagram_nodes) and needs_blank_for_photos and bullets:
            photo_bottom_emu_max = 0
            for p in photos:
                bbox = p.get("bbox")
                if not bbox or len(bbox) != 4:
                    continue
                try:
                    x_norm, y_norm, w_norm, h_norm = bbox
                except Exception:
                    continue

                # Only consider photos that start in the left half; the body box lives on the left.
                if int(x_norm) >= 500:
                    continue

                try:
                    _l, t, _w, h = _bbox_norm_to_emu(
                        [int(x_norm), int(y_norm), int(w_norm), int(h_norm)],
                        slide_w,
                        slide_h,
                    )
                except Exception:
                    continue
                photo_bottom_emu_max = max(photo_bottom_emu_max, int(t + h))

            default_body_top = Inches(2.8)
            body_top = default_body_top
            if photo_bottom_emu_max:
                body_top = max(default_body_top, photo_bottom_emu_max + Inches(0.25))

            footer_gap = Inches(1.0)
            body_height = max(Inches(1.0), int(slide_h - body_top - footer_gap))

            body_box = slide.shapes.add_textbox(
                Inches(0.75),
                body_top,
                Inches(5.4),
                body_height,
            )
            body_box.text_frame.text = bullets
            style_text_frame(
                body_box.text_frame,
                font_color=WHITE,
                margin_top=Inches(0.15),
                auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
                word_wrap=True,
            )

        # Photos: crop from source render and place
        render_path = source_renders_dir / f"slide_{slide_index:02d}.png"
        if photos and render_path.exists():
            with Image.open(render_path) as img:
                for i, photo in enumerate(photos, start=1):
                    try:
                        label = str(photo.get("label") or "").lower()
                        bbox = photo.get("bbox")
                        if not bbox or len(bbox) != 4:
                            continue
                        # Avoid re-introducing full-slide backgrounds; rely on the Premier template backgrounds.
                        if "background" in label or _bbox_area_ratio(bbox) >= 0.85:
                            continue
                        cropped = _crop_bbox_px(img, bbox)
                        crop_path = crops_dir / f"slide_{slide_index:02d}_photo_{i:02d}.png"
                        cropped.save(crop_path, format="PNG")
                        left, top, width, height = _bbox_norm_to_emu(bbox, slide_w, slide_h)
                        slide.shapes.add_picture(str(crop_path), left, top, width=width, height=height)
                    except Exception:
                        continue

        # Diagram nodes as editable shapes (optional)
        title_bottom = 0
        if title_shape_for_bottom is not None:
            try:
                title_bottom = int(title_shape_for_bottom.top + title_shape_for_bottom.height + Inches(0.15))
            except Exception:
                title_bottom = 0
        if use_diagram_nodes:
            _add_diagram_nodes(slide, spec, slide_w, slide_h, WHITE, min_top_emu=title_bottom)

        remove_unused_placeholders(slide)

    # Remove prototype slides from the component library
    for _ in range(original_slide_count):
        delete_slide(prs, 0)

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_pptx))


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--component-library",
        default=str(RUN_ROOT / "assets" / "Premier_Components.pptx"),
        help="Premier component library PPTX",
    )
    parser.add_argument(
        "--specs-dir",
        default=str(RUN_ROOT / "exports" / "conversion" / "slide_specs"),
        help="Directory containing slide_XX.json specs",
    )
    parser.add_argument(
        "--source-renders-dir",
        default=str(RUN_ROOT / "exports" / "conversion" / "source_renders"),
        help="Directory containing slide_XX.png source renders",
    )
    parser.add_argument(
        "--output-pptx",
        default=str(RUN_ROOT / "exports" / "conversion" / "converted_premier_v1.pptx"),
        help="Output converted PPTX",
    )
    parser.add_argument(
        "--crops-dir",
        default=str(RUN_ROOT / "exports" / "conversion" / "crops"),
        help="Where to store cropped photo assets",
    )
    args = parser.parse_args()

    build_deck(
        component_library=Path(args.component_library).resolve(),
        specs_dir=Path(args.specs_dir).resolve(),
        source_renders_dir=Path(args.source_renders_dir).resolve(),
        output_pptx=Path(args.output_pptx).resolve(),
        crops_dir=Path(args.crops_dir).resolve(),
    )
    print(f"Wrote converted deck -> {args.output_pptx}")


if __name__ == "__main__":
    main()
