import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Configuration
RUN_ID = "2026-01-06__accelerating_technology_delivery_presentation"
COMPONENT_LIBRARY = f"runs/{RUN_ID}/assets/Premier_Components.pptx"
OUTPUT_FILE = f"runs/{RUN_ID}/exports/draft_v4_content.pptx"

def delete_slide(prs, index):
    """Delete a slide from the presentation by index."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[index])

def remove_unused_placeholders(slide):
    """Remove placeholder shapes that have no text."""
    # Collect shapes to remove first to avoid modifying list while iterating
    to_remove = []
    for shape in slide.placeholders:
        # Check if it has a text frame and is empty
        if shape.has_text_frame:
            if not shape.text or not shape.text.strip():
                to_remove.append(shape)
        # Keep charts/pictures (logic for another day, currently we only use text placeholders)
    
    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

def style_text_frame(text_frame, font_color=None, margin_left=None, margin_top=None):
    """Apply styling to a text frame."""
    if margin_left is not None:
        text_frame.margin_left = margin_left
    if margin_top is not None:
        text_frame.margin_top = margin_top
    
    # Apply font color to all paragraphs (and future text)
    if font_color:
        for p in text_frame.paragraphs:
            p.font.color.rgb = font_color
            for run in p.runs:
                run.font.color.rgb = font_color

def generate_presentation_v4():
    if not os.path.exists(COMPONENT_LIBRARY):
        print("Error: Component Library not found.")
        return

    # Load the Library
    prs = Presentation(COMPONENT_LIBRARY)
    original_slide_count = len(prs.slides)
    
    # Define Layouts from the Library
    LAYOUT_TITLE = prs.slides[0].slide_layout
    LAYOUT_SECTION = prs.slides[1].slide_layout
    # USE DARK LAYOUTS to match WHITE TEXT
    LAYOUT_CONTENT = prs.slides[5].slide_layout # Dark 1a
    LAYOUT_TWO_COL = prs.slides[13].slide_layout # Dark 1b
    LAYOUT_TWO_COL_SUBTI = prs.slides[9].slide_layout # Dark Subtitle 1a

    print("Generating presentation v4 (Patched Content)...")
    
    WHITE = RGBColor(255, 255, 255)

    # 1. Title Slide (Slide 22)
    # Fix: Clipping (add margin), Contrast (set white), Unused placeholders
    slide = prs.slides.add_slide(LAYOUT_TITLE)
    try:
        title = slide.placeholders[0]
        title.text = "Accelerating Technology Delivery"
        # Fix Title Clipping
        style_text_frame(title.text_frame, font_color=WHITE, margin_left=Inches(0.5))
        
        # Subtitle
        if len(slide.placeholders) > 1:
            sub = slide.placeholders[10]
            sub.text = "Velocity, Quality, and Safety"
            # Fixing clipping with margin
            style_text_frame(sub.text_frame, font_color=WHITE, margin_left=Inches(0.5))
        
        remove_unused_placeholders(slide)
    except Exception as e:
        print(f"Warning Slide 1: {e}")

    # 2. Agenda
    # Fix: Padding/Alignment?
    slide = prs.slides.add_slide(LAYOUT_CONTENT)
    try:
        title = slide.placeholders[0]
        title.text = "Agenda"
        # Since layouts are "Dark", ensure Title is White
        style_text_frame(title.text_frame, font_color=WHITE)

        content = slide.placeholders[10]
        content.text = (
            "1. Current Velocity Metrics\n"
            "2. The Bottleneck Analysis\n"
            "3. Solution: The 'Consensus Loop'\n"
            "4. Roadmap Q1-Q2"
        )
        style_text_frame(content.text_frame, font_color=WHITE, margin_top=Inches(0.2))
        remove_unused_placeholders(slide)
    except Exception as e:
        print(f"Warning Slide 2: {e}")

    # 3. Section 1
    slide = prs.slides.add_slide(LAYOUT_SECTION)
    try:
        # Check placeholder mapping for Section Layout
        # inspect_components says Slide 1 has ph[0] and ph[10]
        title = slide.placeholders[0]
        try:
            # layout might use {{SECTION_TITLE}}
            title.text = "Phase 1: Analysis"
        except:
            # Fallback if mapped incorrectly
            pass 
        style_text_frame(title.text_frame, font_color=WHITE)
        remove_unused_placeholders(slide)
    except Exception as e:
        print(f"Warning Slide 3: {e}")

    # 4. Problem Detail
    slide = prs.slides.add_slide(LAYOUT_CONTENT)
    try:
        slide.placeholders[0].text = "Bottlenecks Identified"
        style_text_frame(slide.placeholders[0].text_frame, font_color=WHITE)
        
        content = slide.placeholders[10]
        content.text = (
            "• Manual Code Reviews: 48h delay average\n"
            "• Flaky Integration Tests: 20% failure rate blocks merge\n"
            "• Monolithic Pipeline: 1h build time slows feedback"
        )
        style_text_frame(content.text_frame, font_color=WHITE, margin_top=Inches(0.2))
        remove_unused_placeholders(slide)
    except Exception as e:
        print(f"Warning Slide 4: {e}")

    # 5. Section 2
    slide = prs.slides.add_slide(LAYOUT_SECTION)
    try:
        slide.placeholders[0].text = "Phase 2: Solution"
        style_text_frame(slide.placeholders[0].text_frame, font_color=WHITE)
        remove_unused_placeholders(slide)
    except Exception as e:
        print(f"Warning Slide 5: {e}")

    # 6. Solution Detail (Two Col)
    slide = prs.slides.add_slide(LAYOUT_TWO_COL)
    try:
        slide.placeholders[0].text = "The Consensus Loop Approach"
        style_text_frame(slide.placeholders[0].text_frame, font_color=WHITE)
        
        # Left
        left = slide.placeholders[1]
        left.text = (
            "Safe-Fail Architecture:\n"
            "• Automated Geometry Linting\n"
            "• Truthful Native Rendering\n"
            "• Visual Critique (AI)"
        )
        style_text_frame(left.text_frame, font_color=WHITE, margin_top=Inches(0.2))
        
        # Right
        right = slide.placeholders[2]
        right.text = (
            "Key Benefits:\n"
            "• 10x Velocity Increase\n"
            "• Zero Brand Defects\n"
            "• Reduced Human Review Time"
        )
        style_text_frame(right.text_frame, font_color=WHITE, margin_top=Inches(0.2))
        
        remove_unused_placeholders(slide)
    except Exception as e:
        print(f"Warning Slide 6: {e}")

    # 7. Section 3
    slide = prs.slides.add_slide(LAYOUT_SECTION)
    try:
        slide.placeholders[0].text = "Phase 3: Execution"
        style_text_frame(slide.placeholders[0].text_frame, font_color=WHITE)
        remove_unused_placeholders(slide)
    except Exception as e:
        print(f"Warning Slide 7: {e}")

    # 8. Roadmap (Two Col Setup)
    # Fix: Padding/Crowding
    slide = prs.slides.add_slide(LAYOUT_TWO_COL_SUBTI)
    try:
        slide.placeholders[0].text = "Q1-Q2 Roadmap"
        style_text_frame(slide.placeholders[0].text_frame, font_color=WHITE)
        
        sub = slide.placeholders[10]
        sub.text = "Scaling the Platform" # Subtitle
        style_text_frame(sub.text_frame, font_color=WHITE)

        # Left
        left = slide.placeholders[1]
        left.text = (
            "Q1: MVP Launch\n"
            "• Deploy Consensus Loop\n"
            "• Onboard Pilot Team\n"
            "• Validate Metrics"
        )
        # Adding Padding to avoid crowding edges
        style_text_frame(left.text_frame, font_color=WHITE, margin_top=Inches(0.2), margin_left=Inches(0.1))

        # Right
        right = slide.placeholders[2]
        right.text = (
            "Q2: Enterprise Scale\n"
            "• Self-Service Portal\n"
            "• Support for 50+ Teams\n"
            "• Advanced Style Governance"
        )
        style_text_frame(right.text_frame, font_color=WHITE, margin_top=Inches(0.2), margin_left=Inches(0.1))
        
        remove_unused_placeholders(slide)
    except Exception as e:
        print(f"Warning Slide 8: {e}")

    # --- CLEANUP: Remove Original Template Slides ---
    print(f"Removing {original_slide_count} template slides...")
    # Delete from 0 to original_count-1.
    # Since deleting index 0 shifts everything down, we just delete index 0 N times.
    for _ in range(original_slide_count):
        delete_slide(prs, 0)

    # Save
    os.makedirs(os.path.dirname(OUTPUT_FILE), exist_ok=True)
    prs.save(OUTPUT_FILE)
    print(f"Draft V4 (Patched) generated at {OUTPUT_FILE}")

if __name__ == "__main__":
    generate_presentation_v4()
