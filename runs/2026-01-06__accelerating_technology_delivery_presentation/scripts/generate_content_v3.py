import os
import sys
from pptx import Presentation

# Configuration
RUN_ID = "2026-01-06__accelerating_technology_delivery_presentation"
COMPONENT_LIBRARY = f"runs/{RUN_ID}/assets/Premier_Components.pptx"
OUTPUT_FILE = f"runs/{RUN_ID}/exports/draft_v3_content.pptx"

def generate_presentation_v3():
    if not os.path.exists(COMPONENT_LIBRARY):
        print("Error: Component Library not found.")
        return

    # Load the Library
    prs = Presentation(COMPONENT_LIBRARY)
    
    # Define Layouts from the Library (mapping mapped via inspect_components.py)
    # Slide 0: Premier Title Dark
    LAYOUT_TITLE = prs.slides[0].slide_layout
    # Slide 1: Premier Title Slide 1 - Dark (Good for Sections)
    LAYOUT_SECTION = prs.slides[1].slide_layout
    # Slide 4: Premier Main Content 1a (Title + 1 Content)
    LAYOUT_CONTENT = prs.slides[4].slide_layout
    # Slide 12: Premier Two Content 1b (Title + 2 Content)
    LAYOUT_TWO_COL = prs.slides[12].slide_layout
    # Slide 8: Premier Two Content Subtitle 1a (Title + Subtitle + 2 Content)
    LAYOUT_TWO_COL_SUBTI = prs.slides[8].slide_layout

    print("Generating presentation v3 (Content Inject)...")

    # The Component Library slides remain in the deck 
    # (or we could start a fresh deck and only use the layouts, but python-pptx 
    # likes keeping the master relations intact from the source).
    # We will append new slides.

    # 1. Title Slide
    slide = prs.slides.add_slide(LAYOUT_TITLE)
    try:
        slide.placeholders[0].text = "Accelerating Technology Delivery"
        # Try finding subtitle placeholder (usually idx 10 or 1)
        # From inspection: Slide 0 has placeholders 0, 10, 11, 12, 13
        if len(slide.placeholders) > 1:
            slide.placeholders[10].text = "Velocity, Quality, and Safety"
    except Exception as e:
        print(f"Warning Slide 1: {e}")

    # 2. Agenda
    slide = prs.slides.add_slide(LAYOUT_CONTENT)
    try:
        slide.placeholders[0].text = "Agenda"
        slide.placeholders[10].text = (
            "1. Current Velocity Metrics\n"
            "2. The Bottleneck Analysis\n"
            "3. Solution: The 'Consensus Loop'\n"
            "4. Roadmap Q1-Q2"
        )
    except Exception as e:
        print(f"Warning Slide 2: {e}")

    # 3. Section 1
    slide = prs.slides.add_slide(LAYOUT_SECTION)
    try:
        slide.placeholders[0].text = "Phase 1: Analysis"
    except Exception as e:
        print(f"Warning Slide 3: {e}")

    # 4. Problem Detail
    slide = prs.slides.add_slide(LAYOUT_CONTENT)
    try:
        slide.placeholders[0].text = "Bottlenecks Identified"
        slide.placeholders[10].text = (
            "• Manual Code Reviews: 48h delay average\n"
            "• Flaky Integration Tests: 20% failure rate blocks merge\n"
            "• Monolithic Pipeline: 1h build time slows feedback"
        )
    except Exception as e:
        print(f"Warning Slide 4: {e}")

    # 5. Section 2
    slide = prs.slides.add_slide(LAYOUT_SECTION)
    try:
        slide.placeholders[0].text = "Phase 2: Solution"
    except Exception as e:
        print(f"Warning Slide 5: {e}")

    # 6. Solution Detail (Two Col)
    slide = prs.slides.add_slide(LAYOUT_TWO_COL)
    try:
        slide.placeholders[0].text = "The Consensus Loop Approach"
        # Left
        slide.placeholders[1].text = (
            "Safe-Fail Architecture:\n"
            "• Automated Geometry Linting\n"
            "• Truthful Native Rendering\n"
            "• Visual Critique (AI)"
        )
        # Right
        slide.placeholders[2].text = (
            "Key Benefits:\n"
            "• 10x Velocity Increase\n"
            "• Zero Brand Defects\n"
            "• Reduced Human Review Time"
        )
    except Exception as e:
        print(f"Warning Slide 6: {e}")

    # 7. Section 3
    slide = prs.slides.add_slide(LAYOUT_SECTION)
    try:
        slide.placeholders[0].text = "Phase 3: Execution"
    except Exception as e:
        print(f"Warning Slide 7: {e}")

    # 8. Roadmap (Two Col Setup)
    slide = prs.slides.add_slide(LAYOUT_TWO_COL_SUBTI)
    try:
        slide.placeholders[0].text = "Q1-Q2 Roadmap"
        slide.placeholders[10].text = "Scaling the Platform" # Subtitle
        # Left
        slide.placeholders[1].text = (
            "Q1: MVP Launch\n"
            "• Deploy Consensus Loop\n"
            "• Onboard Pilot Team\n"
            "• Validate Metrics"
        )
        # Right
        slide.placeholders[2].text = (
            "Q2: Enterprise Scale\n"
            "• Self-Service Portal\n"
            "• Support for 50+ Teams\n"
            "• Advanced Style Governance"
        )
    except Exception as e:
        print(f"Warning Slide 8: {e}")

    # Save
    os.makedirs(os.path.dirname(OUTPUT_FILE), exist_ok=True)
    prs.save(OUTPUT_FILE)
    print(f"Draft V3 generated at {OUTPUT_FILE}")

if __name__ == "__main__":
    generate_presentation_v3()
