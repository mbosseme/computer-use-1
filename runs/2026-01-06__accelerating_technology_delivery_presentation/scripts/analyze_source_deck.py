from __future__ import annotations

import argparse
import hashlib
import json
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


RUN_ROOT = Path(__file__).resolve().parents[1]


@dataclass
class ImageRef:
    slide_index: int
    shape_id: int
    filename: str
    sha256: str
    content_type: str
    width_emu: int
    height_emu: int
    left_emu: int
    top_emu: int


@dataclass
class ShapeSummary:
    shape_id: int
    shape_type: str
    name: str
    left_emu: int
    top_emu: int
    width_emu: int
    height_emu: int
    has_text: bool
    text_preview: str


@dataclass
class SlideSummary:
    slide_index: int
    layout_name: str
    width_emu: int
    height_emu: int
    shape_count: int
    picture_count: int
    text_shape_count: int
    approx_picture_area_ratio: float
    looks_like_flat_image_slide: bool
    shapes: List[ShapeSummary]


def _sha256(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _safe_text_preview(text: str, limit: int = 120) -> str:
    text = (text or "").strip().replace("\n", " ")
    if len(text) <= limit:
        return text
    return text[: limit - 1] + "…"


def _shape_type_name(shape) -> str:
    try:
        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.PICTURE:
            return "PICTURE"
        if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return "AUTO_SHAPE"
        if st == MSO_SHAPE_TYPE.TEXT_BOX:
            return "TEXT_BOX"
        if st == MSO_SHAPE_TYPE.PLACEHOLDER:
            return "PLACEHOLDER"
        if st == MSO_SHAPE_TYPE.GROUP:
            return "GROUP"
        if st == MSO_SHAPE_TYPE.TABLE:
            return "TABLE"
        if st == MSO_SHAPE_TYPE.CHART:
            return "CHART"
        return str(st)
    except Exception:
        return "UNKNOWN"


def _shape_geometry(shape) -> Tuple[int, int, int, int]:
    left = int(getattr(shape, "left", 0) or 0)
    top = int(getattr(shape, "top", 0) or 0)
    width = int(getattr(shape, "width", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)
    return left, top, width, height


def _shape_has_text(shape) -> bool:
    try:
        return bool(getattr(shape, "has_text_frame", False)) and bool(shape.text_frame)
    except Exception:
        return False


def _get_shape_text(shape) -> str:
    if not _shape_has_text(shape):
        return ""
    try:
        return shape.text_frame.text or ""
    except Exception:
        return ""


def _picture_area_ratio(slide, slide_w: int, slide_h: int) -> float:
    if slide_w <= 0 or slide_h <= 0:
        return 0.0
    slide_area = float(slide_w * slide_h)
    total = 0.0
    for shp in slide.shapes:
        if shp.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        _, _, w, h = _shape_geometry(shp)
        total += float(max(w, 0) * max(h, 0))
    return min(1.0, total / slide_area)


def _looks_like_flat_image_slide(
    slide,
    slide_w: int,
    slide_h: int,
    picture_area_ratio: float,
    picture_count: int,
    text_shape_count: int,
) -> bool:
    # Heuristic: one big picture covering most of the slide, with little/no text.
    if picture_count == 0:
        return False
    if picture_count == 1 and picture_area_ratio >= 0.75 and text_shape_count <= 1:
        return True
    if picture_area_ratio >= 0.85 and text_shape_count == 0:
        return True
    return False


def extract_images(prs: Presentation, out_dir: Path) -> List[ImageRef]:
    out_dir.mkdir(parents=True, exist_ok=True)
    image_refs: List[ImageRef] = []
    seen: Dict[str, str] = {}

    for slide_index, slide in enumerate(prs.slides, start=1):
        for shp in slide.shapes:
            if shp.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                img = shp.image
                blob = img.blob
                digest = _sha256(blob)
                ext = (img.ext or "bin").lower()
                filename = seen.get(digest)
                if filename is None:
                    filename = f"img_{digest[:12]}.{ext}"
                    (out_dir / filename).write_bytes(blob)
                    seen[digest] = filename

                left, top, w, h = _shape_geometry(shp)
                image_refs.append(
                    ImageRef(
                        slide_index=slide_index,
                        shape_id=int(getattr(shp, "shape_id", 0) or 0),
                        filename=filename,
                        sha256=digest,
                        content_type=getattr(img, "content_type", "") or "",
                        width_emu=w,
                        height_emu=h,
                        left_emu=left,
                        top_emu=top,
                    )
                )
            except Exception:
                # Some picture-like shapes may not expose image data; skip gracefully.
                continue

    return image_refs


def summarize(prs: Presentation) -> List[SlideSummary]:
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)

    summaries: List[SlideSummary] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        shapes: List[ShapeSummary] = []
        picture_count = 0
        text_shape_count = 0

        for shp in slide.shapes:
            left, top, w, h = _shape_geometry(shp)
            text = _get_shape_text(shp)
            has_text = bool(text.strip())
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_count += 1
            if has_text:
                text_shape_count += 1

            shapes.append(
                ShapeSummary(
                    shape_id=int(getattr(shp, "shape_id", 0) or 0),
                    shape_type=_shape_type_name(shp),
                    name=str(getattr(shp, "name", "") or ""),
                    left_emu=left,
                    top_emu=top,
                    width_emu=w,
                    height_emu=h,
                    has_text=has_text,
                    text_preview=_safe_text_preview(text),
                )
            )

        picture_area_ratio = _picture_area_ratio(slide, slide_w, slide_h)
        looks_flat = _looks_like_flat_image_slide(
            slide,
            slide_w,
            slide_h,
            picture_area_ratio=picture_area_ratio,
            picture_count=picture_count,
            text_shape_count=text_shape_count,
        )

        summaries.append(
            SlideSummary(
                slide_index=slide_index,
                layout_name=str(getattr(slide.slide_layout, "name", "") or ""),
                width_emu=slide_w,
                height_emu=slide_h,
                shape_count=len(slide.shapes),
                picture_count=picture_count,
                text_shape_count=text_shape_count,
                approx_picture_area_ratio=picture_area_ratio,
                looks_like_flat_image_slide=looks_flat,
                shapes=shapes,
            )
        )

    return summaries


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--input",
        default=str(RUN_ROOT / "inputs" / "presentation_to_convert" / "Accelerating Technology Delivery.pptx"),
        help="Path to the PPTX to analyze",
    )
    parser.add_argument(
        "--out-json",
        default=str(RUN_ROOT / "exports" / "conversion" / "source_analysis.json"),
        help="Where to write the analysis JSON",
    )
    parser.add_argument(
        "--extract-images-dir",
        default=str(RUN_ROOT / "exports" / "conversion" / "source_images"),
        help="Directory to extract embedded images",
    )
    args = parser.parse_args()

    input_path = Path(args.input).resolve()
    out_json = Path(args.out_json).resolve()
    images_dir = Path(args.extract_images_dir).resolve()
    out_json.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(input_path))

    slide_summaries = summarize(prs)
    image_refs = extract_images(prs, images_dir)

    payload: Dict[str, Any] = {
        "input": str(input_path),
        "slide_count": len(prs.slides),
        "slide_width_emu": int(prs.slide_width),
        "slide_height_emu": int(prs.slide_height),
        "slides": [asdict(s) for s in slide_summaries],
        "extracted_images_dir": str(images_dir),
        "images": [asdict(i) for i in image_refs],
    }
    out_json.write_text(json.dumps(payload, indent=2), encoding="utf-8")

    flat = [s.slide_index for s in slide_summaries if s.looks_like_flat_image_slide]
    print(f"Analyzed {len(prs.slides)} slides")
    print(f"Extracted {len(image_refs)} picture occurrences -> {images_dir}")
    if flat:
        print(f"Slides that look like flattened image slides (heuristic): {flat}")
    else:
        print("No slides flagged as flattened image slides by heuristic")
    print(f"Wrote analysis JSON -> {out_json}")


if __name__ == "__main__":
    main()
